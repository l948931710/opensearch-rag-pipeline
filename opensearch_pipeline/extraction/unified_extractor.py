# -*- coding: utf-8 -*-
"""
unified_extractor.py — 统一文档提取入口

单一入口 UnifiedExtractor.extract()，根据 file_ext 分发到对应的提取器。
支持 mock 模式（接收 mock_text）和生产模式（读取真实文件）。

DAG 层不需要知道文件类型，只调用 extract() 即可。
"""

import copy
import datetime as _dt
import hashlib
import os
import re
import tempfile
from typing import Dict, List, Optional

from opensearch_pipeline.extraction.schema import ExtractionResult, ExtractedBlock
from opensearch_pipeline.extraction.text_extractor import (
    extract_text_file,
    blocks_to_text,
    extract_title_from_blocks,
)
from opensearch_pipeline.extraction.ocr_client import OCRClient, sanitize_ocr_text
from opensearch_pipeline import funnel_verdict_store
from opensearch_pipeline.ingest_flags import (
    effective_funnel_policy_version,
    pdf_strip_stitch_enabled,
)


_DEFAULT_IMG_DIR: Optional[str] = None


def _safe_image_output_dir(task: dict) -> str:
    """解析图片导出目录。绝不返回 ""——空串会被 os.path.join("", filename) 落到 cwd,
    污染工作区(历史上往仓库根目录散落 *_img*.png / *_slide*.png)。调用方未传 `_tmp_dir`
    时,回退到一个稳定的系统临时目录(进程级复用,不落仓库)。"""
    d = (task or {}).get("_tmp_dir")
    if d:
        return d
    global _DEFAULT_IMG_DIR
    if _DEFAULT_IMG_DIR is None:
        _DEFAULT_IMG_DIR = tempfile.mkdtemp(prefix="rag_extract_images_")
    return _DEFAULT_IMG_DIR


# ── 基础抽取器依赖表：这三类文档的【文本抽取】100% 依赖对应第三方库 ──
# 缺失时 _extract_xlsx/_extract_pptx/_extract_docx 的兜底 except 会把 ImportError
# 吞成 "Failed to extract ...: No module named '...'" warning，产出 0 块空 canonical
# 且 stage-1 照常全绿——环境性失败被优雅降级掩盖（2026-07-16 现场：/usr/bin/python3
# 缺 openpyxl）。preflight_extractor_deps 供 stage-1 抽取节点在处理任何文档前
# fail-fast。PIL/OCR/图片导出等辅助依赖有意不在此列：辅助失败不破坏文本抽取
# （graceful degradation 铁律）。pdf 也有意不预检：pdfplumber→pypdf 双库互为降级、
# 之后还有 OCR 兜底，单库缺失不构成总损；双库全缺产出的 "pypdf/PyPDF2 not installed"
# 空 canonical 由 node_build_canonical 的 ENV-DEP 守卫兜住。
BASE_EXTRACTOR_DEPS = {
    "xlsx": ("openpyxl", "openpyxl"),
    "pptx": ("pptx", "python-pptx"),
    "docx": ("docx", "python-docx"),
}


def preflight_extractor_deps(file_exts, _import=None) -> List[tuple]:
    """对将要处理的文件类型预检基础抽取器依赖。

    Args:
        file_exts: 本批任务涉及的扩展名集合（大小写/前导点均容忍）。
        _import: 测试注入点，默认 importlib.import_module。
    Returns:
        缺失清单 [(ext, module_name, pip_name), ...]（按 ext 排序）；全可用时为 []。
        只查 import 可用性，不做版本校验；fail-fast 还是告警由调用方决定。
    """
    import importlib
    imp = _import or importlib.import_module
    missing = []
    for ext in sorted({(e or "").lower().strip().lstrip(".") for e in file_exts}):
        dep = BASE_EXTRACTOR_DEPS.get(ext)
        if not dep:
            continue
        module_name, pip_name = dep
        try:
            imp(module_name)
        except ImportError:
            missing.append((ext, module_name, pip_name))
    return missing


_PUA_RANGE = (0xE000, 0xF8FF)          # Private Use Area：坏字体常映射到此
_CID_RE = re.compile(r"\(cid:\d+\)")   # pdfminer 对无 ToUnicode 映射字体的输出


def _native_text_quality_ok(text: str) -> bool:
    """G15：原生文本质量谓词（纯字符串数学，零 API）。

    旧门槛只看字符数（≥50 即"有文本"）：坏内嵌字体页吐出的 (cid:123) 串 /
    U+FFFD / PUA 乱码 ≥50 字符就永不重 OCR，垃圾直接进切块与向量。
    判 False 的三个独立信号（任一命中即视为乱码页）：
      1. (cid:N) 模式覆盖原文 >15%；
      2. U+FFFD 替换符 + PUA 字符密度 >10%；
      3. CJK/ASCII可打印/中文标点/全半角 合计占比 <0.70。
    """
    stripped = "".join(text.split())
    if not stripped:
        return True  # 空文本交给字符数门槛处理，此谓词不重复触发
    cid_chars = sum(len(m.group(0)) for m in _CID_RE.finditer(text))
    if cid_chars / max(len(text), 1) > 0.15:
        return False
    total = len(stripped)
    junk = good = 0
    for ch in stripped:
        cp = ord(ch)
        if cp == 0xFFFD or _PUA_RANGE[0] <= cp <= _PUA_RANGE[1]:
            junk += 1
            continue
        if (0x20 <= cp <= 0x7E                      # ASCII 可打印
                or 0x4E00 <= cp <= 0x9FFF           # CJK 统一表意
                or 0x3400 <= cp <= 0x4DBF           # CJK 扩展 A
                or 0x3000 <= cp <= 0x303F           # CJK 标点
                or 0xFF00 <= cp <= 0xFFEF           # 全半角形式
                or 0x2010 <= cp <= 0x2027           # 常用横线/引号/省略号区
                or ch in "℃±×÷·—…" ):
            good += 1
    if junk / total > 0.10:
        return False
    return good / total >= 0.70


def _detect_password_protected(local_path: str, file_ext: str) -> Optional[str]:
    """G14：加密/带密码文件前置探测（纯本地字节嗅探 + pypdf is_encrypted，零 API 成本）。

    此前加密文件只表现为 pdfplumber/pypdf/python-docx 的通用打开失败：确定性失败
    仍吃满 ≤3 次重试，加密 PDF 还会被 _pages_needing_ocr 的 [1] 占位空烧一次 OCR。
    返回 None = 未加密或探测失败（fail-open，走原抽取路径）；返回字符串 = 加密原因。
    """
    if not local_path or not os.path.exists(local_path):
        return None
    try:
        if file_ext == "pdf":
            from pypdf import PdfReader
            reader = PdfReader(local_path)
            if not reader.is_encrypted:
                return None
            # 空密码可解（仅 owner-password 限权、内容可读）→ 放行走正常抽取
            try:
                if int(reader.decrypt("")) != 0:
                    return None
            except Exception:
                pass
            return "PDF 受用户密码保护，无法解密读取"
        if file_ext in ("docx", "xlsx", "pptx"):
            # OOXML 本质是 ZIP（PK 头）；Office 加密后被包进 OLE 复合文档
            # （EncryptedPackage 流），文件头变为 D0 CF 11 E0 A1 B1 1A E1。
            with open(local_path, "rb") as f:
                head = f.read(8)
            if head == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
                return (f"{file_ext} 文件头为 OLE 复合文档——Office 加密文档"
                        f"（EncryptedPackage）或旧版二进制格式（doc/xls/ppt）误改后缀")
            return None
    except Exception:
        return None  # 探测自身失败不拦路（fail-open）
    return None


def _iter_pptx_shapes(shapes):
    """G13：深度优先展开 GroupShape（组合形状）。

    此前主循环只遍历 slide.shapes 顶层：GroupShape 无 .text，被
    `hasattr(shape, "text")` 静默跳过——组内文本框（中文 SOP/培训 deck 里
    "框+箭头+说明文字"的主流写法）整体丢失。组内嵌套组递归展开，子形状
    按 XML 文档序原位替换组本身，表格/占位符判定与顶层同语义。"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        try:
            is_group = shape.shape_type == MSO_SHAPE_TYPE.GROUP
        except Exception:
            is_group = False
        if is_group:
            yield from _iter_pptx_shapes(shape.shapes)
        else:
            yield shape


def _warn_skipped_vector_images(warnings: list, stats: dict) -> None:
    """G4：矢量图（WMF/EMF/SVG）被跳过时落一条结构化 warning，使内容损失可见。

    此前三处裸 continue 静默丢弃——中文企业 DOCX 里剪贴板粘贴的 ERP 截图/Visio
    流程图常为 EMF，丢了完全无感知。计数进 result.warnings（随 canonical 落库、
    可被治理看板/审计聚合）；跳过行为本身不变（栅格化转换是后续工作）。"""
    n = int(stats.get("skipped_vector_images", 0) or 0)
    if n <= 0:
        return
    formats = "/".join(sorted(stats.get("skipped_vector_formats", ()) or ())) or "vector"
    warnings.append(
        f"[VECTOR_IMAGE_SKIPPED] {n} 张矢量图（{formats}）未提取——"
        f"WMF/EMF/SVG 暂不支持栅格化，图内内容未进入索引。"
    )


def _warn_page_cap_skips(warnings: list, stats: dict) -> list:
    """B5（2026-07-25）：图片提取的页上限跳过 → 结构化 warning + partial_loss_notes 追加项。

    与 OCR 页上限、原生抽取页上限走**同一条**留痕通道：三处此前都只打日志/只写结构化字段，
    没有一处进 partial_loss_notes ⇒ 长文档尾部内容不进索引，而 document_version 照常
    DONE+INDEXED，运维无法用一条 SQL 圈出需要重扫的对象。
    """
    n = int((stats or {}).get("skipped_pages_beyond_cap", 0) or 0)
    if n <= 0:
        return []
    note = (f"image_page_cap: {n} page(s) beyond pdf_image_max_pages were never scanned "
            f"for images — figures on those pages are missing from the index")
    warnings.append(note)
    return [note]


def _warn_pdf_image_extract_failed(warnings: list, img_stats: dict) -> list:
    """审查 P1-8：PDF 图片提取**整段失败**（缺 PyMuPDF / 打不开文件）→ 结构化 warning +
    partial_loss_notes 追加项。

    不留这条的话，「一张图都没提取出来」与「这篇本来就没图」在 canonical 里长得一模一样
    （都是 `assets: []`），下游的资产集比对会把前者读成"基线的图全没了"。
    """
    reason = (img_stats or {}).get("pdf_image_extract_failed")
    if not reason:
        return []
    note = (f"pdf_image_extract_failed: {reason} —— 本篇**一张图都未提取**，"
            f"assets 为空不代表原文无图")
    warnings.append(note)
    return [note]


def _warn_verdict_store_unavailable(warnings: list, funnel_stats: dict) -> list:
    """E：权威判决记录不可用且无同策略缓存兜底 → warning + partial_loss_notes 追加项。

    为什么必须显式降级而不是静默继续（codex BLOCKER-3）：这一格里本次的图集是**现场
    重判**的结果，可能与该文档已服务版本不一致，而"靠方案 F 事后喊出来"不成立——
    F 依赖同一个 RDS、自身 fail-open，且正文改一个字就只产 source_changed 普通审计。
    落 NEEDS_REVIEW 让人来看，与 ocr_partial / vlm_degraded 同一通道。
    """
    note = (funnel_stats or {}).get("verdict_store_note")
    if not note:
        return []
    warnings.append(note)
    return [note]


def _warn_unreadable_images(warnings: list, funnel_stats: dict) -> list:
    """A11（2026-07-25）：把「解码失败」的图记成结构化 warning，并返回 partial_loss_notes 追加项。

    此前这类图被 `_static_heuristics` 的 `(0,0,0.0)` 兜底压成 DISCARD_DECORATIVE，
    与"真的是装饰线条"混为一谈：无 warning、无 partial note、不计 vlm_degraded，
    运维只能在日志里看到一行和正常丢弃一模一样的 print。

    ⚠️ 文案里的 `decode failed` 是**载荷**不是修辞：node_write_chunk_meta 的 0-chunk
    「疑似失败」判据按关键词表（fail/error/cannot/无法/错误/失败）扫 canonical warnings，
    命中才落 chunk_status='NEEDS_REVIEW' + content_process_status='FAILED'。改文案前先看
    那份关键词表（tests 有钉死这条耦合）。
    """
    names = list(funnel_stats.get("unreadable_names") or [])
    if not names:
        return []
    shown = "; ".join(names[:5]) + (f" ...(+{len(names) - 5})" if len(names) > 5 else "")
    note = (f"[IMAGE_UNREADABLE] {len(names)} 张图片无法解码（decode failed），"
            f"未进入漏斗/索引：{shown}")
    warnings.append(note)
    return [note]


def _html_to_text(raw_html: str) -> str:
    """HTML → 纯文本：剥标签、跳过 script/style、块级标签转换行；失败回退原文。"""
    import re
    from html.parser import HTMLParser

    class _TextExtractor(HTMLParser):
        _BLOCK_TAGS = {
            "p", "div", "br", "li", "tr", "table", "section", "article",
            "ul", "ol", "h1", "h2", "h3", "h4", "h5", "h6",
        }

        def __init__(self):
            super().__init__(convert_charrefs=True)
            self.parts = []
            self._skip_depth = 0

        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style"):
                self._skip_depth += 1
            elif tag in self._BLOCK_TAGS:
                self.parts.append("\n")

        def handle_endtag(self, tag):
            if tag in ("script", "style"):
                self._skip_depth = max(0, self._skip_depth - 1)
            elif tag in self._BLOCK_TAGS:
                self.parts.append("\n")

        def handle_data(self, data):
            if not self._skip_depth:
                self.parts.append(data)

    try:
        parser = _TextExtractor()
        parser.feed(raw_html)
        parser.close()
        text = "\n".join(ln.strip() for ln in "".join(parser.parts).splitlines())
        text = re.sub(r"\n{3,}", "\n\n", text).strip()
        return text or raw_html
    except Exception:
        return raw_html


def _csv_to_text(raw_csv: str) -> str:
    """CSV → 行式文本：csv 模块处理引号/转义，单元格以 " | " 连接；失败回退原文。"""
    import csv
    import io

    try:
        try:
            dialect = csv.Sniffer().sniff(raw_csv[:4096], delimiters=",;\t")
        except csv.Error:
            dialect = csv.excel
        rows = csv.reader(io.StringIO(raw_csv), dialect)
        lines = [
            " | ".join(cell.strip() for cell in row)
            for row in rows if any((cell or "").strip() for cell in row)
        ]
        return "\n".join(lines) or raw_csv
    except Exception:
        return raw_csv


def _read_text_with_fallback(local_path: str):
    """读文本文件（txt/md/csv/html）的编码回退链，返回 (text, warnings)。

    utf-8-sig 严格 → gb18030 严格 → utf-8 errors=replace + 可见警告。
    此前单读 utf-8 errors="ignore"：GBK/GB2312（中文 Windows/Excel 导出默认编码）
    的 CJK 字节几乎全被静默丢弃 → 文档被误判为空、或以乱码 DONE 入索引。
    gb18030 是 GBK/GB2312 的超集且对纯 ASCII 无歧义（utf-8 先试保证真 utf-8 不会
    被误判）；最后一档 replace 保留可解部分并打警告，让文档被标记而非静默糊掉。
    """
    with open(local_path, "rb") as f:
        data = f.read()
    warnings: List[str] = []
    for enc in ("utf-8-sig", "gb18030"):
        try:
            text = data.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    else:
        text = data.decode("utf-8", errors="replace")
        warnings.append(
            "Text decode fallback: file is neither valid utf-8 nor gb18030; "
            "decoded with utf-8 errors=replace, content may be corrupted"
        )
    # 二进制读 + decode 不做 universal-newline 转换，自补（下游按 \n 分块）
    return text.replace("\r\n", "\n").replace("\r", "\n"), warnings


# 缓存条目「可消费」状态：process_image 的四值契约（image_funnel_processor.py 的
# 返回契约 docstring）。刻意**不含** DISCARD_UNREADABLE —— 写入门本就拒它（解码失败
# 多为一次性），万一历史脏写混进缓存，重算才是正确处置。
_VLM_CACHE_USABLE_STATUSES = {"DISCARD_DECORATIVE", "ROUTE_TO_TEXT",
                              "ROUTE_TO_VECTOR", "QUARANTINE_SENSITIVE"}

# public 裸 key 回退**专用**白名单：刻意排除 QUARANTINE_SENSITIVE（裸 key 全部产生于
# public-bypass 时代，复用它会跳过敏感审计）。与上面那个集合语义不同，**禁止合并**。
_VLM_CACHE_LEGACY_PUBLIC_STATUSES = {"DISCARD_DECORATIVE", "ROUTE_TO_TEXT", "ROUTE_TO_VECTOR"}


def _vlm_cache_entry_usable(entry) -> bool:
    """缓存条目可消费性契约（单一来源）。

    D1（2026-07-25）：消费者直接吃 entry.get("ocr_text") 与 entry["status"]，条目形态
    一旦异常就不是"这张图没缓存"，而是**整篇文档抽取抛异常失败**（实测：本机缓存里
    3 条 `{sha256}:pub:2` 的值是 JSON 数字，命中即 AttributeError）。此处把形态校验收
    到唯一入口，非法条目一律视同未命中 → 重算 → 同 key 被合法 dict 覆盖，自愈。

    校验到 status 而不止 isinstance(dict)：未知状态不会抛异常，但会安静地绕过
    DISCARD 等值跳过、带着空字段流进 asset_dict —— 那比崩溃更难查。

    先判 str 再判集合成员：status 若是 list/dict 等不可哈希值，直接 `in <set>` 会抛
    TypeError —— 守卫自己崩掉就等于没修（本函数的单测钉死这条）。
    """
    if not isinstance(entry, dict):
        return False
    status = entry.get("status")
    return isinstance(status, str) and status in _VLM_CACHE_USABLE_STATUSES


def _vlm_cache_ns(is_public):
    """缓存命名空间：pub/sec [+ 版本标签]。

    升级 VLM 模型或改 funnel prompt 后，把 RAG_VLM_CACHE_VERSION 设成新值即可让旧缓存
    整体【干净失效】、按新模型重打，而无需手动清 scratch/OSS 缓存。
    B4（生产级外审 2026-07-17 P2-03）：默认版本 "2"——配合 file_hash MD5→SHA-256 切换
    让存量条目整体失效（构造 MD5 碰撞继承他图 CLEAN 结论的通道关闭）。代价=存量图下次
    被触碰时重过 VLM 审计（增量、按再摄取面，非一次性全量）。设 "" 可还原历史裸键形。
    """
    ns = "pub" if is_public else "sec"
    ver = os.environ.get("RAG_VLM_CACHE_VERSION", "2").strip()
    ns = f"{ns}:{ver}" if ver else ns
    # 漏斗策略后缀（选项 C）。⚠️ **2026-07-26 起 C 默认 ON ⇒ 默认键就带 `:c1`**
    # （此处旧注释曾写"默认策略为空串、存量键逐字节不变"，那是 C 默认 OFF 时的事实，
    # 翻默认后与代码相反 —— 审查遗漏项已订正）。存量键因此不再被主查询命中，靠
    # `_cross_policy_cache_fallback` 有条件复用。
    # 键必须带策略：命中判定发生在 process_image 之前，不带的话开了 C 也会被旧的 DISCARD
    # 条目直接短路 —— A/B 变成假 ON，反向（C-ON 写的 rescue 被 C-OFF 读到）同样不成立。
    policy = effective_funnel_policy_version()
    return f"{ns}:{policy}" if policy else ns


def _funnel_doc_title(task: dict) -> str:
    """VLM funnel 的文档语境（#F-mm7，RAG_VLM_DOC_CONTEXT 默认 OFF）。

    修死参：extraction task 构造从不设 doc_title 键 → funnel prompt 的【文档标题】
    上下文分支从未生效，visual_summary/image_category 缺文档语境，拉低内容匹配绑定
    （xlsx P0 IDF-bigram / 图N bigram）与 VL rerank 依赖的图注质量。工厂文件名
    （如「注塑机安全操作规程.docx」）就是最强的免费文档语境。

    OFF（默认）→ 沿用 task.get('doc_title','') = 恒空串，行为逐字节一致。
    ⚠️ 开启前置（eval-first，缺一不放行）：
      1. 同步设新 RAG_VLM_CACHE_VERSION（VLM 输入变化 → 旧缓存 caption 必须干净失效）；
      2. 跑 L4-ingestion 全格式闸：xlsx ≥0.8917 / docx ≥0.95 硬线（caption 变化会
         移动内容匹配层的绑定结果——xlsx 0.8917 是按现状 caption 调优的）；
      3. 已知偏置（显式接受）：VLM 缓存 key=(md5, ns) 不含文档语境，同 MD5 图片跨
         文档复用时 caption 由先处理的文档标题决定（first-wins），换 version 清不掉
         ——prompt 侧已加「不要照抄标题」指令压低串染面。
    """
    explicit = task.get("doc_title") or ""
    if os.environ.get("RAG_VLM_DOC_CONTEXT", "").strip().lower() not in ("1", "true", "yes"):
        return explicit   # OFF：与历史 task.get("doc_title", "") 完全一致（恒 ""）
    if explicit:
        return explicit
    fn = os.path.basename(task.get("filename") or "")
    return os.path.splitext(fn)[0]


def _xlsx_cell_to_str(c) -> str:
    """xlsx 单元格值 → 字符串。data_only=True 下日期单元格是 datetime/date 对象，裸 str() 会产出
    '2024-01-15 00:00:00' 这类噪声 → 渲染为干净日期。
    注意：百分比/前导零(001234) 的语义依赖 cell.number_format，而 read_only 的 iter_rows(values_only)
    路径取不到 number_format（大表强制 read_only）→ 那部分不在此处理（B2-4 deferred，避免只在小表生效
    的不一致 + 改 chunk 文本带来的 embedding 漂移）。仅日期清理：不依赖 number_format，两条路径一致、安全。"""
    if isinstance(c, _dt.datetime):
        if (c.hour, c.minute, c.second) == (0, 0, 0):
            return c.strftime("%Y-%m-%d")
        return c.strftime("%Y-%m-%d %H:%M")
    if isinstance(c, _dt.date):
        return c.strftime("%Y-%m-%d")
    return str(c)


def _image_sha256(local_path: str, fallback_seed=None) -> str:
    """图片字节 sha256（VLM 缓存主键）。读不出来时回退到进程内唯一键 —— 该键以
    `fallback_` 开头，写缓存的门会据此拒绝入库（跨进程无意义的键不能落盘）。"""
    try:
        with open(local_path, "rb") as f:
            return hashlib.sha256(f.read()).hexdigest()
    except Exception:
        return "fallback_%s" % (id(fallback_seed) if fallback_seed is not None else id(local_path))


def _vlm_cache_put_allowed(simulate: bool, funnel_res: dict, file_hash: str) -> bool:
    """VLM 缓存写入门（**单一来源** —— 嵌入图批量路径与独立图片路径共用）。

    拒绝四类：
      · simulate 的 mock 结果（按真实 sha256 落键后会毒化真实运行）；
      · degraded（VLM 超时/解析失败的兜底）—— 缓存会把一次性故障固化成永久标签；
      · A11 的 DISCARD_UNREADABLE —— 同上，解码失败往往是一次性的；
      · fallback_ 前缀键（基于内存地址，跨进程无意义）。
    """
    return (not simulate
            and not funnel_res.get("degraded")
            and funnel_res.get("status") != "DISCARD_UNREADABLE"
            and not str(file_hash).startswith("fallback_"))


def _vlm_cache_entry(funnel_res: dict) -> dict:
    """缓存条目的字段形态（单一来源，防两条路径各写各的）。"""
    return {
        "status": funnel_res.get("status", ""),
        "visual_summary": funnel_res.get("visual_summary", ""),
        "image_category": funnel_res.get("image_category", ""),
        "vlm_annotation_map": funnel_res.get("vlm_annotation_map", {}),
        "reason": funnel_res.get("reason", ""),
        "width": funnel_res.get("width", 0),
        "height": funnel_res.get("height", 0),
        "file_size_kb": funnel_res.get("file_size_kb", 0.0),
        "ocr_text": funnel_res.get("ocr_text", ""),
    }


def _funnel3_discard(entry) -> bool:
    """该缓存条目是否是 **Funnel 3 的语义弃图** —— 也就是选项 C 会改变结论的那一类。

    非 DISCARD 一律 False（C 不动 ROUTE_TO_* / QUARANTINE_*）。是 DISCARD 但
    **认不出 Funnel 级别**时返回 True（保守重判）：宁可多付一次 VLM，也不能让一条
    来路不明的旧弃图在新策略下静默沿用 —— 错弃是不可见且不可恢复的那一侧。
    """
    if not isinstance(entry, dict):
        return False
    if not str(entry.get("status", "")).startswith("DISCARD"):
        return False
    # Funnel 1 的弃图是纯几何判据，确定且与 C 无关，可安全跨策略复用
    return "Funnel 1" not in str(entry.get("reason", ""))


def _cross_policy_cache_fallback(vlm_cache, file_hash, is_public):
    """非默认漏斗策略下，回退查基础策略键 —— **仅当 C 不可能改变该条目结论时**接受。

    命中即迁移到当前策略键（内存中；真实运行随 _save_vlm_cache 持久化）。
    默认策略下本函数是 no-op（当前键就是基础键）。
    """
    policy = effective_funnel_policy_version()
    if not policy:
        return None
    ns = "pub" if is_public else "sec"
    ver = os.environ.get("RAG_VLM_CACHE_VERSION", "2").strip()
    base_key = f"{file_hash}:{ns}:{ver}" if ver else f"{file_hash}:{ns}"
    entry = vlm_cache.get(base_key)
    if not _vlm_cache_entry_usable(entry) or _funnel3_discard(entry):
        return None
    vlm_cache[f"{file_hash}:{_vlm_cache_ns(is_public)}"] = entry
    return entry


def _vlm_cache_lookup(vlm_cache, file_hash, is_public):
    """带命名空间的 VLM 缓存查询 + 跨策略/遗留裸 MD5 key 的条件回退。

    遗留条目（命名空间化之前写入，OSS 上 ~1500 条）用现行带后缀 key 永远
    查不到 → 回灌时全量重打 VLM 白花钱。回退仅限 public 命名空间（裸 key
    全部产生于 public-bypass 时代；quarantine 的 :sec 绝不回退，避免跳过
    敏感审计），且要求 status 合法、非 QUARANTINE_SENSITIVE、无 Simulated
    污染。命中即迁移到带后缀 key（内存中；真实运行随 _save_vlm_cache
    持久化，simulate 不落盘），裸 key 自然老化。只读新写仍只用带后缀 key。
    """
    cache_key = f"{file_hash}:{_vlm_cache_ns(is_public)}"
    entry = vlm_cache.get(cache_key)
    if entry is not None:
        if _vlm_cache_entry_usable(entry):
            return entry
        # D1：形态异常的条目**视同未命中**（不 early-return，继续走下面的回退逻辑
        # = 忠实地"就当这个 key 不存在"）。advisory 缓存绝不能打断抽取。
        # 只打类型与截断后的 status，不输出 caption/OCR 正文。
        _bad = entry.get("status") if isinstance(entry, dict) else None
        print(f"      ⚠️ [VLM Cache] 条目形态异常，视同未命中重算: {cache_key[:20]}… "
              f"(type={type(entry).__name__}, status={str(_bad)[:24]!r})")

    # ── 跨策略条件回退（选项 C）──
    # 整体换策略命名空间的代价是"5 万条缓存一次性失效 → 下次摄取全量重付 VLM，且把每张图
    # 重新暴露给今天那个已被证明更差的 VLM 制度"。而现在挡住漂移的**恰恰是** 2026-07-20
    # 那批缓存条目（方案 §8.3）—— 整体换 key 等于亲手把保护拆掉。
    # 所以只对"C 可能改变结论"的条目重判：Funnel-3 的 DISCARD。Funnel-1 的结构性弃图是
    # 确定的几何判据、C 不碰；ROUTE_TO_* / QUARANTINE_* 也不受 C 影响，可安全复用。
    # 判不出是哪一级的 DISCARD 一律**保守重判**（宁可多花一次调用，绝不静默沿用旧弃图）。
    policy_entry = _cross_policy_cache_fallback(vlm_cache, file_hash, is_public)
    if policy_entry is not None:
        return policy_entry

    # 设了显式缓存版本（模型/prompt 升级）→ 不回退任何旧 key，让旧条目干净失效、按新模型重打。
    if os.environ.get("RAG_VLM_CACHE_VERSION", "").strip():
        return None
    if not is_public:
        return None
    legacy = vlm_cache.get(file_hash)
    if not isinstance(legacy, dict):
        return None
    if legacy.get("status") not in _VLM_CACHE_LEGACY_PUBLIC_STATUSES:
        return None
    blob = (f"{legacy.get('visual_summary', '')}|{legacy.get('ocr_text', '')}"
            f"|{legacy.get('reason', '')}")
    if "simulat" in blob.lower():
        return None
    vlm_cache[cache_key] = legacy
    return legacy


def _funnel1_discard_bound(w, h, kb):
    """Funnel-1 静态丢弃判据（w<50 / h<50 / <3KB / aspect>8）的单点定义。

    与 _process_embedded_images 的 Funnel-1 过滤及 image_funnel_processor 的
    heuristic_prefilter 语义一致；DOCX（_stitch_strip_runs）与 PDF
    （_stitch_pdf_strip_assets）两个缝合器共用 —— 「逐条都是今日 Funnel-1 丢弃
    对象」是缝合的 load-bearing 闸（只救回本来必死的图，绝不动能存活的图），
    分头手写第三份会漂移。
    """
    aspect = max(w / max(h, 1), h / max(w, 1))
    return w < 50 or h < 50 or kb < 3.0 or aspect > 8.0


def _stitch_pdf_strip_assets(assets):
    """缝合 PDF 条带切片（#F-mm9，`RAG_PDF_STRIP_STITCH`，**2026-07-26 起默认 ON**）
    → (assets, warnings)。

    翻默认的依据（方案 §15-16，Sam 拍板 2026-07-26）：
      · FL-XS-WI-007 实测：33 条横条缝成 **2 张完整的《交货单》**（手册号/客户/商检号/货号/
        表格/条形码俱全），存活 8→10 零丢失。而《交货单》正是该篇《吸塑扫码报检》SOP 的
        核心物件（作业说明/步骤2/4/5.1/6 逐条点名）—— 今天它是**静默全损**的；
      · L4（GT 修正后）：该篇 0.7778 → 0.8889，pdf 总 0.89236 → 0.91320，其余 5 篇逐位不变，
        `img_dup` 不动；
      · **生产影响面 = 恰好 1 篇**（580 篇当前版本 PDF 全扫，实扫 563 篇，命中率 0.18%，
        命中的就是这一篇本身）；
      · 结构上只碰 `_funnel1_discard_bound` 已判死的资产 —— **不可能伤害任何当前存活的图**。

    回滚：显式设 `0/false/no/off`。

    PDF 转换器把一张照片存成 N 条 1000×31 全宽横条时，每条 aspect>8 逐条死于
    Funnel-1，整图彻底消失且零告警（实证：FL-XS-WI-007 PDF 的 23 条横条）。
    DOCX 版（_stitch_strip_runs）靠 blocks 连续 image_ref run 判连续性；PDF 的
    ref 是 chunk 期才注入、无 blocks 信号，但资产带 page_num + 显示 bbox ——
    改用【同页 + bbox 纵向邻接】判据（下一条 y0≈上一条 y1，坐标信号比 DOCX 的
    连续性推断更强），插在 extract_images_from_pdf 之后、_process_embedded_images
    之前（资产级，缝合产物照常过漏斗/VLM/上传/绑定）。

    判据（移植 DOCX 版四闸，全部满足才缝合）：
      - 同页、bbox 纵向邻接（tol 3pt）且 x 区间对齐（纵向堆叠的前提）；
      - run 内任一条 bbox 缺失（get_image_info 可失败）→ 该 run 整体放弃；
      - 像素宽完全一致；逐条 height<80 且条状（w≥3h 或 w≥100）；
      - 逐条都是 Funnel-1 丢弃对象（_funnel1_discard_bound）——只救必死的图；
      - 缝合总高 [100, 6000]。
    ⚠️ md5 去重资产洞（已知交互，显式接受）：extract_images_from_pdf 按
    (md5, 朝向op) 全局去重发生在本函数之前，条带中字节相同的条（纯色带）已被
    丢弃 → 邻接断裂即视为 run 边界，一张图可能缝成多个各自完整的子段（碎片化
    子缝合优于整图消失），断裂子 run 仍须各自满足全部判据。
    缝合产物：bbox=成员并集（正常参与 chunk 期几何 y 锚定）、保留 page_num 与
    首条 image_index。任何 PIL/IO 失败该 run 原样保留（fail-open）。
    """
    # 2026-07-26 起**默认 ON**（Sam 拍板）。off-list 极性：默认 ON 的开关若把未知值也当关，
    # 配错一个字就会静默退回"整图消失且零告警"，而这条链本来就无信号可查。
    if not pdf_strip_stitch_enabled():
        return assets, []
    if not assets:
        return assets, []
    try:
        from PIL import Image
    except ImportError:
        return assets, []

    ADJ_TOL = 3.0        # bbox 纵向邻接 / x 对齐容差（pt）
    MIN_RUN = 4
    MAX_SLICE_H = 80

    # 逐资产读像素尺寸/大小（失败的资产不参与缝合、原样保留）
    info = {}   # id(asset) -> (w, h, kb)
    for a in assets:
        lp = getattr(a, "local_path", "")
        if not lp or not os.path.exists(lp):
            continue
        try:
            with Image.open(lp) as im:
                w, h = im.size
            info[id(a)] = (w, h, os.path.getsize(lp) / 1024.0)
        except Exception:
            continue

    # 按页分组 → 页内按 bbox y0 排序 → 邻接切 run
    by_page = {}
    for a in assets:
        by_page.setdefault(getattr(a, "page_num", None), []).append(a)

    def _eligible(a):
        """可参与 run 的最低门槛：bbox 齐备 + 像素信息可读 + 逐条判据全过。"""
        if getattr(a, "bbox", None) is None or id(a) not in info:
            return False
        w, h, kb = info[id(a)]
        return (h < MAX_SLICE_H and (w >= 3 * h or w >= 100)
                and _funnel1_discard_bound(w, h, kb))

    runs = []
    for page, page_assets in by_page.items():
        if page is None:
            continue
        cands = sorted((a for a in page_assets if _eligible(a)),
                       key=lambda a: (a.bbox[1], a.bbox[0]))
        cur = []
        for a in cands:
            if not cur:
                cur = [a]
                continue
            prev = cur[-1]
            first = cur[0]
            same_width = info[id(a)][0] == info[id(first)][0]
            adjacent = abs(a.bbox[1] - prev.bbox[3]) <= ADJ_TOL
            x_aligned = (abs(a.bbox[0] - first.bbox[0]) <= ADJ_TOL
                         and abs(a.bbox[2] - first.bbox[2]) <= ADJ_TOL)
            if same_width and adjacent and x_aligned:
                cur.append(a)
            else:
                if len(cur) >= MIN_RUN:
                    runs.append(cur)
                cur = [a]
        if len(cur) >= MIN_RUN:
            runs.append(cur)

    if not runs:
        return assets, []

    warnings = []
    drop_ids = set()
    replacement = {}   # id(first_asset) -> composite asset
    for run in runs:
        total_h = sum(info[id(a)][1] for a in run)
        if not (100 <= total_h <= 6000):
            continue
        w0 = info[id(run[0])][0]
        try:
            canvas = Image.new("RGB", (w0, total_h), (255, 255, 255))
            y = 0
            for a in run:
                with Image.open(a.local_path) as im:
                    canvas.paste(im.convert("RGB"), (0, y))
                y += info[id(a)][1]
            first = run[0]
            out_path = f"{os.path.splitext(first.local_path)[0]}_pdfstitched{len(run)}.png"
            canvas.save(out_path)
        except Exception as e:
            print(f"      ⚠️ [pdf-strip-stitch] PIL stitch failed (run kept as-is): {e}")
            continue
        comp = copy.copy(first)
        comp.local_path = out_path
        comp.original_name = f"stitched:{len(run)}slices"
        comp.bbox = (
            min(a.bbox[0] for a in run), min(a.bbox[1] for a in run),
            max(a.bbox[2] for a in run), max(a.bbox[3] for a in run),
        )
        replacement[id(first)] = comp
        for a in run:
            drop_ids.add(id(a))
        msg = (f"[pdf-strip-stitch] p{first.page_num}: {len(run)} slices "
               f"({w0}x{total_h}) → {os.path.basename(out_path)}")
        warnings.append(msg)
        print(f"      {msg}")

    if not replacement:
        return assets, []

    out = []
    for a in assets:
        if id(a) in replacement:
            out.append(replacement[id(a)])   # 首条位置放缝合产物（保持文档序）
        elif id(a) not in drop_ids:
            out.append(a)
    return out, warnings


def _stitch_strip_runs(blocks, image_assets, min_run=4, max_slice_height=80):
    """缝合 Word 条带切片：一张照片被存成 N 条同宽窄条时，逐条都会被
    Funnel-1 丢弃（h<50 或 aspect>8），整图彻底丢失。

    判据（每条都 load-bearing，全部满足才缝合，否则原样返回）：
      - blocks 中 ≥min_run 个连续 image_ref（之间无任何其他块——图片独占段落
        不产生文本块，所以"连续"是唯一可观测信号）；
      - 资产像素宽完全一致（纵向堆叠的前提）；
      - 逐条 height<max_slice_height 且条状（w≥3h 或 w≥100）——方形图标行
        （工具栏/①②③枚举）不会被误缝；
      - 逐条都是今日 Funnel-1 的丢弃对象（w<50/h<50/<3KB/aspect>8）——
        只"救回"，绝不改变现状能存活的图；
      - 缝合总高在 [100, 6000] 内。
    网格切片（如 23×31 小方块）因不满足条状判据而不缝——纵向堆叠会产生
    无意义的细长柱，留待网格重建（future work）。
    任何 PIL/IO 失败 → 该 run 原样保留（fail-open）。
    """
    if not blocks or not image_assets:
        return blocks, image_assets
    try:
        from PIL import Image
    except ImportError:
        return blocks, image_assets

    def _btype(b):
        return b.block_type if hasattr(b, "block_type") else b.get("block_type", "")

    def _bextra(b):
        e = b.extra if hasattr(b, "extra") else b.get("extra")
        return e if isinstance(e, dict) else {}

    asset_by_index = {}
    for a in image_assets:
        asset_by_index.setdefault(getattr(a, "image_index", None), a)

    runs = []  # [run_blocks, ...]
    cur = []
    for b in blocks:
        if _btype(b) == "image_ref":
            cur.append(b)
        else:
            if len(cur) >= min_run:
                runs.append(list(cur))
            cur = []
    if len(cur) >= min_run:
        runs.append(list(cur))
    if not runs:
        return blocks, image_assets

    _discard_bound = _funnel1_discard_bound

    drop_block_ids = set()
    drop_asset_idx = set()
    new_assets = []

    for run_blocks in runs:
        infos = []  # (block, asset, w, h, kb, image_index)
        ok = True
        for b in run_blocks:
            idx = _bextra(b).get("image_index")
            a = asset_by_index.get(idx)
            lp = getattr(a, "local_path", "") if a else ""
            if not a or not lp or not os.path.exists(lp):
                ok = False
                break
            try:
                with Image.open(lp) as im:
                    w, h = im.size
                kb = os.path.getsize(lp) / 1024.0
            except Exception:
                ok = False
                break
            infos.append((b, a, w, h, kb, idx))
        if not ok or not infos:
            continue

        if len({w for _, _, w, _, _, _ in infos}) != 1:
            continue
        if not all(
            h < max_slice_height and (w >= 3 * h or w >= 100) and _discard_bound(w, h, kb)
            for _, _, w, h, kb, _ in infos
        ):
            continue
        w0 = infos[0][2]
        total_h = sum(h for _, _, _, h, _, _ in infos)
        if not (100 <= total_h <= 6000):
            continue

        try:
            canvas = Image.new("RGB", (w0, total_h), (255, 255, 255))
            y = 0
            for _, a, _, h, _, _ in infos:
                with Image.open(a.local_path) as im:
                    canvas.paste(im.convert("RGB"), (0, y))
                y += h
            first_a = infos[0][1]
            out_path = f"{os.path.splitext(first_a.local_path)[0]}_stitched{len(infos)}.png"
            canvas.save(out_path)
        except Exception as e:
            print(f"      ⚠️ [strip-stitch] PIL stitch failed (run kept as-is): {e}")
            continue

        first_b = infos[0][0]
        first_idx = infos[0][5]
        for _, _, _, _, _, idx in infos:
            drop_asset_idx.add(idx)
        comp = copy.copy(first_a)
        comp.local_path = out_path
        comp.image_index = first_idx
        comp.original_name = f"stitched:{len(infos)}slices"
        new_assets.append(comp)
        extra0 = first_b.extra if hasattr(first_b, "extra") else None
        if isinstance(extra0, dict):
            extra0["stitched_from"] = len(infos)
        for b in run_blocks[1:]:
            drop_block_ids.add(id(b))
        print(f"      [strip-stitch] {len(infos)} slices ({w0}x{total_h}) → "
              f"{os.path.basename(out_path)}")

    if not new_assets:
        return blocks, image_assets

    blocks = [b for b in blocks if id(b) not in drop_block_ids]
    image_assets = [
        a for a in image_assets
        if getattr(a, "image_index", None) not in drop_asset_idx
    ] + new_assets
    return blocks, image_assets


# 双重 OCR 判重阈值：funnel ocr_text 块的字符 3-gram 有 ≥85% 出现在同页页级 OCR
# 文本里即视为"内容已被整页转写覆盖"。实测（6B37DC 重复对）0.987 / 无关文本 0.0，
# 两侧余量都大；OCR 双通道字符级差异（换行/全半角标点）被 3-gram+空白归一吸收。
_DOUBLE_OCR_CONTAINMENT_THRESHOLD = 0.85


def _drop_double_ocr_dups(blocks: list, page_ocr_blocks: list) -> tuple:
    """整页扫描 PDF 双重 OCR 去重（G6 tier-3 决策抓获，docs/tier3_decision_2026-07-06.md）。

    纯扫描页会走两条 OCR 路径：①页面作为嵌入图片被 image funnel ROUTE_TO_TEXT
    转写一次（blocks 里带 extra.source_image 的 ocr_text 块）②per-page OCR 兜底又
    对同一页整页转写一次——两份文本只差换行/标点（字节 hash 不同，精确去重不吃），
    canonical 内容翻倍（生产实锤 6B37DC：1 页 600 字 → 11 chunk/1501 字，"第七条"×4）。

    以页级 OCR 为该页的权威整页转写（整页渲染必然覆盖页内嵌入图），剔除同页 funnel
    产物中内容已被其覆盖的 ocr_text 块。只剔文本块——asset 本身不动（图片照常
    上传/绑定/在答案里渲染）。fail-open：文本过短无指纹（gram_containment 返 0）/
    页级 OCR 未覆盖该页 / 内容确有差异 → 保留 funnel 块（宁重勿丢，遵循 OCR/funnel
    失败绝不阻断文本抽取的约定）。

    Args:
        blocks: 当前 result.blocks（含 funnel 产物）。
        page_ocr_blocks: per-page OCR 兜底新产出的块（page_num 齐全）。

    Returns:
        (kept_blocks, dropped_pages)：过滤后的 blocks 与被剔除 funnel 块的页码列表。
    """
    from opensearch_pipeline.text_similarity import gram_containment

    page_ocr_text: Dict[int, List[str]] = {}
    for b in page_ocr_blocks:
        pg = getattr(b, "page_num", None)
        if pg is not None:
            page_ocr_text.setdefault(pg, []).append(getattr(b, "text", "") or "")
    if not page_ocr_text:
        return blocks, []

    kept, dropped_pages = [], []
    for b in blocks:
        extra = getattr(b, "extra", None) or {}
        pg = getattr(b, "page_num", None)
        if (getattr(b, "block_type", "") == "ocr_text"
                and extra.get("source_image")     # funnel ROUTE_TO_TEXT 产物特征
                and pg in page_ocr_text):
            try:
                containment = gram_containment(
                    getattr(b, "text", "") or "", "\n".join(page_ocr_text[pg]))
            except Exception:
                containment = 0.0  # 判重原语异常 → fail-open 保留
            if containment >= _DOUBLE_OCR_CONTAINMENT_THRESHOLD:
                dropped_pages.append(pg)
                continue
        kept.append(b)
    return kept, dropped_pages


_HF_DIGIT_NORM = re.compile(r"\d+")


def _strip_hf_lines_from_ocr(ocr_blocks: list, hf_texts: set) -> tuple:
    """G12-OCR：把原生 Pass-1 检出的页眉/页脚词集应用到页级 OCR 文本（G6 tier-3 决策抓获②）。

    原生 pdfplumber 路径按 y 边界裁页眉页脚，OCR 兜底文本**无坐标**——整页信头
    （公司名/文件编号/版本/生效日期）被逐字转进 canonical（实测 doc2-p2 页眉字符
    ≈5 倍于正文）。本函数用 Pass-1 的跨页词集（词级 + 数字归一，"FL-CW-SW-003"→
    "FL-CW-SW-###"）做行级过滤：某行的全部多字符 token（同款归一后）都 ∈ 词集
    即视为页眉/页脚行剔除（单字符 token 如 年/月/日/A// 不参与判定，随行走）。

    fail-open 三道保险：①词集为空（纯扫描件原生零文本 → Pass-1 无词集）= no-op；
    ②词集规模阀：>60 项视为页眉检测器失灵（真实信头是紧凑的），整体放弃过滤——
    刻意【不用】按块行数比例阀：迷你页（如尾页只有一句话）的 OCR 本来就页眉占
    多数，比例阀会杀死主修复场景；③过滤后整块为空才移除块（正文块不受影响）。
    返回 (blocks, dropped_line_count)。
    """
    if not hf_texts or len(hf_texts) > 60:
        return ocr_blocks, 0
    dropped_total = 0
    out = []
    for b in ocr_blocks:
        text = getattr(b, "text", "") or ""
        lines = text.splitlines()
        keep, dropped = [], 0
        for ln in lines:
            toks = [t for t in ln.split() if len(t) > 1]
            if toks and all(_HF_DIGIT_NORM.sub("#", t) in hf_texts for t in toks):
                dropped += 1
                continue
            keep.append(ln)
        if dropped:
            dropped_total += dropped
            new_text = "\n".join(keep).strip()
            if not new_text:
                continue           # 整块皆页眉 → 移除块
            b.text = new_text
        out.append(b)
    return out, dropped_total


class UnifiedExtractor:
    """
    统一文档提取器。

    用法：
        extractor = UnifiedExtractor(simulate=True)
        result = extractor.extract(task)

    task dict 需包含：
        doc_id, version_no, file_ext, raw_key
        mock_text (可选，模拟模式)
        local_path (可选，生产模式)
    """

    # 原生提取文本低于此阈值 → OCR fallback（图片类型用整体阈值）
    OCR_THRESHOLD_CHARS = 100
    # Increment 0b: PDF 逐页 OCR 门槛——原生文本少于此字符数的页视为扫描页/坏字体页，
    # 单独送 OCR；取代旧的"按整文档 text_length 一刀切"（封面有字则整文档跳过 OCR）。
    PER_PAGE_OCR_THRESHOLD = 50
    # 原生文本抽取的页上限（G2 env 化：RAG_PDF_NATIVE_MAX_PAGES）。类属性是默认值，
    # __init__ 会以 config 值覆盖为实例属性（类级兜底保护 __new__ 构造的 gate-only 用法）；
    # extract_pdf 调用、截断标注与 OCR 扫描钳位共用同一值。
    # OCR 门槛只考察【原生抽取实际覆盖过】的页：超过此上限的页根本没被抽取（按设计截断），
    # 不能因其"0 原生字符"就误判为需 OCR——那会浪费 Qwen-VL 预算、并把真扫描页挤出 OCR 配额。
    pdf_native_max_pages = 1000
    # PDF 嵌入图片挖掘页上限（RAG_PDF_IMAGE_MAX_PAGES）：产出图逐张进 OCR+VLM 付费漏斗。
    pdf_image_max_pages = 100

    def __init__(
        self,
        oss_client=None,
        ocr_client: Optional[OCRClient] = None,
        simulate: bool = None,
    ):
        from opensearch_pipeline.config import get_config
        cfg = get_config()
        if simulate is None:
            simulate = cfg.simulate
        self.oss_client = oss_client
        if not ocr_client:
            self.ocr_client = OCRClient(
                api_key=cfg.ocr.api_key,
                api_base_url=cfg.ocr.api_base_url,
                ocr_model=cfg.ocr.model,
                max_ocr_pages=cfg.ocr.max_ocr_pages,
                simulate=simulate,
            )
        else:
            self.ocr_client = ocr_client
        self.simulate = simulate
        self.config = cfg
        # G2：PDF 页上限从 config 注入（getattr 兜底保护直接构造 PipelineConfig 的旧测试桩）
        self.pdf_native_max_pages = int(getattr(cfg, "pdf_native_max_pages", 1000) or 1000)
        self.pdf_image_max_pages = int(getattr(cfg, "pdf_image_max_pages", 100) or 100)
        # 可选：运行级成本熔断器（由 orchestrator 注入，跨文档累计运行预算）。
        # 为空时 vlm_rebuilder 会按 cfg 现造一个做单文档闸。
        self.cost_breaker = None

    def extract(self, task: dict) -> ExtractionResult:
        """
        统一提取入口。

        根据 file_ext 分发，自动处理 OCR fallback。
        """
        file_ext = task.get("file_ext", "txt").lower().strip().lstrip(".")

        # ── Mock 模式：直接解析注入的文本 ──
        if "mock_text" in task:
            return self._extract_mock(task, file_ext)

        # ── 生产模式：按文件类型分发 ──
        # G14：加密/带密码文档前置短路——显式 [PASSWORD_PROTECTED] 警告 + 不触发
        # OCR/VLM 付费调用；extract_method 复用 unsupported 前缀使 0-chunk 收尾按
        # 既有关键词落 NEEDS_REVIEW/FAILED 等人工提供解密版本（而非无谓自动重试）。
        if file_ext in ("pdf", "docx", "xlsx", "pptx"):
            _enc_reason = _detect_password_protected(task.get("local_path", ""), file_ext)
            if _enc_reason:
                return ExtractionResult(
                    doc_id=task["doc_id"],
                    version_no=task["version_no"],
                    source_key=task.get("raw_key", ""),
                    file_ext=file_ext,
                    extract_method=f"unsupported:password_protected_{file_ext}",
                    title=task.get("filename", ""),
                    text="",
                    text_length=0,
                    ocr_required=False,
                    ocr_status="NOT_REQUIRED",
                    warnings=[f"[PASSWORD_PROTECTED] {_enc_reason}；请提供解密版本后重新上传"],
                )
        if file_ext == "pdf":
            return self._extract_pdf(task)
        elif file_ext == "docx":
            return self._extract_docx(task)
        elif file_ext == "xlsx":
            return self._extract_xlsx(task)
        elif file_ext == "xls":
            # 旧版二进制 Excel：按用户决策不在管线内支持 —— .xls/.doc 走一次性
            # 转换（xlsx/docx）后回灌；ingest_policy 已把 xls 排除在扫描之外。
            # 显式 unsupported 保证错误可见（勿误路由进 _extract_xlsx 静默吞掉）。
            return self._unsupported(task, file_ext)
        elif file_ext == "pptx":
            return self._extract_pptx(task)
        elif file_ext in ("txt", "md", "csv", "html", "htm"):
            return self._extract_text(task)
        elif file_ext in ("png", "jpg", "jpeg", "webp", "tif", "tiff", "gif", "bmp"):
            return self._extract_image(task)
        else:
            return self._unsupported(task, file_ext)

    # ── Mock 模式 ──

    def _extract_mock(self, task: dict, file_ext: str) -> ExtractionResult:
        """解析 mock_text 为结构化 blocks。"""
        text = task["mock_text"]
        blocks = extract_text_file(text, source="mock")
        title = extract_title_from_blocks(blocks, fallback=task.get("filename", ""))
        flat_text = blocks_to_text(blocks)

        return ExtractionResult(
            doc_id=task["doc_id"],
            version_no=task["version_no"],
            source_key=task.get("raw_key", ""),
            file_ext=file_ext,
            extract_method="mock_injection",
            title=title,
            text=flat_text,
            text_length=len(flat_text),
            blocks=blocks,
            page_count=None,
            ocr_required=False,
            ocr_status="NOT_REQUIRED",
        )

    # ── PDF ──

    def _extract_pdf(self, task: dict) -> ExtractionResult:
        """PDF 提取（文本 + 嵌入图片 + OCR fallback）。"""
        from opensearch_pipeline.extraction.pdf_extractor import extract_pdf
        from opensearch_pipeline.extraction.image_extraction_utils import extract_images_from_pdf

        local_path = task.get("local_path", "")
        _pdf_meta: dict = {}   # G12-OCR：接 Pass-1 页眉/页脚词集（供 OCR 兜底文本裁剪）
        blocks, page_count, warnings = extract_pdf(
            local_path, max_pages=self.pdf_native_max_pages, meta_out=_pdf_meta)
        flat_text = blocks_to_text(blocks)
        title = extract_title_from_blocks(blocks, fallback=task.get("filename", ""))

        # 提取嵌入图片 →（可选）条带缝合 → 三阶段过滤漏斗
        # #F-mm9：缝合必须在漏斗之前（条带逐条会死于 Funnel-1 aspect>8），
        # 缝合产物照常过漏斗/VLM/上传/绑定；默认 OFF 时 _stitch_pdf_strip_assets 原样透传。
        img_stats: dict = {}
        raw_assets = extract_images_from_pdf(
            local_path, _safe_image_output_dir(task), max_pages=self.pdf_image_max_pages,
            stats=img_stats)
        _warn_skipped_vector_images(warnings, img_stats)
        _page_cap_notes = _warn_page_cap_skips(warnings, img_stats)   # B5：图片页上限留痕
        _page_cap_notes.extend(_warn_pdf_image_extract_failed(warnings, img_stats))  # P1-8
        raw_assets, stitch_warnings = _stitch_pdf_strip_assets(raw_assets)
        warnings.extend(stitch_warnings)
        _funnel_stats: dict = {}
        assets, img_blocks, vlm_degraded = self._process_embedded_images(
            raw_assets, task, stats=_funnel_stats)
        _warn_unreadable_images(warnings, _funnel_stats)   # A11：坏图留痕（PDF）
        _page_cap_notes.extend(_warn_verdict_store_unavailable(warnings, _funnel_stats))  # E
        if img_blocks:
            blocks.extend(img_blocks)
            flat_text = blocks_to_text(blocks)

        # 判断实际使用的提取方法
        has_layout_blocks = any(
            b.extra.get("detected_by") in ("font_size", "bold_font", "pdfplumber_lines", "layout")
            for b in blocks if hasattr(b, "extra") and b.extra
        )
        extract_method = "pdfplumber_layout" if has_layout_blocks else "pypdf"

        result = ExtractionResult(
            doc_id=task["doc_id"],
            version_no=task["version_no"],
            source_key=task.get("raw_key", ""),
            file_ext="pdf",
            extract_method=extract_method,
            title=title,
            text=flat_text,
            text_length=len(flat_text),
            blocks=blocks,
            page_count=page_count,
            warnings=warnings,
            assets=assets,
            vlm_degraded_count=vlm_degraded,
        )
        # Stash 本地路径，供 _pages_needing_ocr 在 page_count<=0 时
        # 走保守 fallback 重拿真实页数（不暴露成 dataclass field 以避免序列化噪声）
        result._local_path = local_path  # type: ignore[attr-defined]
        # G12-OCR：页眉/页脚词集随 result 走到 _apply_ocr_fallback（同 _garbled_pages 惯例）
        _hf = (_pdf_meta.get("header_texts") or set()) | (_pdf_meta.get("footer_texts") or set())
        if _hf:
            result._hf_texts = frozenset(_hf)  # type: ignore[attr-defined]


        # OCR fallback 判断（perf#65：needy 单算传递——此前 _needs_ocr / _apply_ocr_fallback
        # 各算一遍，全 blocks 逐页聚合与 page_count<=0 的 pypdf/pdfplumber 回退探测都跑两遍）
        needy_pages = self._pages_needing_ocr(result)
        if needy_pages:
            result = self._apply_ocr_fallback(task, result, needy_pages)

        # ── Increment 1: VLM 版面重建（逐页升级不可提取页）──
        # 默认关闭（cfg.rebuild.enabled=False）→ no-op；开启后受成本熔断器约束。
        try:
            from opensearch_pipeline.extraction.vlm_rebuilder import maybe_rebuild_pdf
            result = maybe_rebuild_pdf(task, result, self.config,
                                       breaker=getattr(self, "cost_breaker", None))
        except Exception as _e:
            print(f"    ⚠️ [vlm_rebuilder] skipped (non-fatal): {_e}", flush=True)

        # ── Increment 2: VLM 表格精修（结构错乱的 PDF 表格；数字保真闸把关）──
        # 默认关闭（cfg.rebuild.refine_tables=False）→ 完全 no-op（零回归）。
        try:
            from opensearch_pipeline.extraction.vlm_rebuilder import maybe_refine_tables
            result = maybe_refine_tables(task, result, self.config,
                                         breaker=getattr(self, "cost_breaker", None))
        except Exception as _e:
            print(f"    ⚠️ [table_refine] skipped (non-fatal): {_e}", flush=True)

        # P1-09：标注原生抽取页上限截断。page_count 是【真实总页数】（extract_pdf 切片前计算），
        # 但原生抽取全失败时 extract_pdf 返回 page_count=0，_pages_needing_ocr 会用 pypdf/pdfplumber
        # 恢复真实页数并回写 result.page_count（入口局部 page_count 仍是 0）——若只看局部量，>20 页
        # 扫描 PDF 会漏判截断被当完整上线。真实总页数优先取 result.page_count，回退入口 page_count，
        # 并回写 result.page_count 确保总页数落库（document_version.page_count），供治理看板统计截断文档数。
        # B5：图片页上限的留痕合并进来（**追加不覆盖** —— 同一 closure 里还有 A11 的
        # unreadable、批次6 的 ocr_partial、B4 的 table_drop）
        if _page_cap_notes:
            result.partial_loss_notes.extend(_page_cap_notes)

        effective_page_count = result.page_count or page_count or 0
        if effective_page_count > self.pdf_native_max_pages:
            result.extract_truncated = True
            result.extracted_pages = self.pdf_native_max_pages
            result.page_count = effective_page_count
            _trunc_note = (
                f"[TRUNCATED] PDF 共 {effective_page_count} 页，仅前 {self.pdf_native_max_pages} 页被抽取；"
                f"第 {self.pdf_native_max_pages + 1} 页起的内容未进入索引（含 OCR）。")
            result.warnings.append(_trunc_note)
            # B5（2026-07-25）：此前只有 warnings + 结构化字段，**没进 partial_loss_notes** ——
            # 于是这类文档照常定稿 DONE，不走 NEEDS_REVIEW 复查通道，运维圈不出重扫对象。
            # 三处页级上限（原生抽取 / OCR / 图片提取）自此统一走同一条留痕通道。
            result.partial_loss_notes.append(
                f"native_page_cap: {effective_page_count - self.pdf_native_max_pages} page(s) "
                f"beyond pdf_native_max_pages were never extracted")

        return result

    # ── DOCX ──

    def _extract_docx(self, task: dict) -> ExtractionResult:
        """DOCX 提取（文本 + 嵌入图片，段落级图片位置追踪）。"""
        from opensearch_pipeline.extraction.docx_extractor import extract_docx_with_images
        from opensearch_pipeline.extraction.image_extraction_utils import extract_images_from_docx

        local_path = task.get("local_path", "")

        # perf#63：docx.Document 全量解析（OOXML 解包 + lxml 树构建）只做一次，注入给
        # 文本位置追踪与图片导出两个函数复用（此前各自 load，一个文档解析两遍）。
        # 依赖缺失/坏文档时传 None → 两函数各自走旧的自行加载路径并产出原有 warning
        # （"Failed to open DOCX" 3-tuple 留痕行为不变，graceful degradation）。
        _document = None
        try:
            import docx as _docx
            _document = _docx.Document(local_path)
        except Exception:
            _document = None

        # ── 方案 B：用 extract_docx_with_images 获得段落级 image_ref 位置 ──
        # 这让 _inject_image_ref_blocks 走精确匹配路径而非启发式均匀分配
        blocks, inline_image_assets, docx_warnings = extract_docx_with_images(
            local_path, document=_document)
        warnings = list(docx_warnings)
        flat_text = blocks_to_text(blocks)
        title = extract_title_from_blocks(blocks, fallback=task.get("filename", ""))

        # ── 提取嵌入图片到磁盘 → 三阶段过滤漏斗 ──
        # perf#63：复用同一 _document，省掉此前这里的第二次全量解析。
        img_stats: dict = {}
        exported_images = extract_images_from_docx(
            local_path, _safe_image_output_dir(task), document=_document, stats=img_stats
        )
        _warn_skipped_vector_images(warnings, img_stats)

        # ── 对齐 image_index：inline_image_assets 按文档顺序，
        #    exported_images 按 rels 遍历顺序，需要通过 target_ref 匹配 ──
        if inline_image_assets and exported_images:
            # 构建 target_ref → exported asset 映射
            export_by_ref = {}
            for ea in exported_images:
                ref = getattr(ea, "original_name", "") or ""
                if ref:
                    export_by_ref[ref] = ea

            # 用 inline 顺序重建 exported 列表，确保 image_index 一致。
            # 同一 media 被正文引用多次时（同 rId 多处出现），第二次起必须
            # 复制资产对象——直接复用会让 image_index 互相覆盖（last-write-wins），
            # 先出现的步骤永远绑不到图。copy.copy 保留 XLSX 等动态附加属性。
            aligned_exports = []
            consumed_ids = set()
            for ia in inline_image_assets:
                ref = getattr(ia, "original_name", "") or ""
                matched = export_by_ref.get(ref)
                if matched:
                    if id(matched) in consumed_ids:
                        matched = copy.copy(matched)
                    else:
                        consumed_ids.add(id(matched))
                    # 用 inline 的 image_index 覆盖 export 的 image_index
                    matched.image_index = ia.image_index
                    aligned_exports.append(matched)

            # 如果对齐成功（大部分都能匹配），用对齐后的列表——但【保留】未匹配到 inline ref 的
            # 导出图（页眉/页脚/VML-only、或 ref 解析失败者），否则它们被静默丢弃、永不索引/渲染。
            # 追加在末尾并重排 image_index（接 aligned 之后），不碰 inline 已绑定图的下标/顺序。
            if len(aligned_exports) >= len(exported_images) * 0.5:
                leftovers = [ea for ea in exported_images if id(ea) not in consumed_ids]
                next_idx = max((getattr(ea, "image_index", 0) for ea in aligned_exports),
                               default=-1) + 1
                for ea in leftovers:
                    ea.image_index = next_idx
                    next_idx += 1
                exported_images = aligned_exports + leftovers

        # ── 条带切片缝合（须在对齐之后：依赖 ref↔asset 的 image_index 对应）──
        try:
            blocks, exported_images = _stitch_strip_runs(blocks, exported_images)
        except Exception as _e:
            print(f"      ⚠️ [strip-stitch] skipped (non-fatal): {_e}")

        _funnel_stats: dict = {}
        assets, img_blocks, vlm_degraded = self._process_embedded_images(
            exported_images, task, stats=_funnel_stats,
        )
        _warn_unreadable_images(warnings, _funnel_stats)   # A11：坏图留痕（DOCX）
        _verdict_notes = _warn_verdict_store_unavailable(warnings, _funnel_stats)   # E
        # 不再 extend img_blocks — image_ref 已内联在 blocks 中
        # img_blocks 是 _process_embedded_images 生成的冗余 image_ref，
        # 只有当 blocks 中完全没有 image_ref 时才 fallback 追加
        has_inline_refs = any(
            (b.block_type if hasattr(b, "block_type") else b.get("block_type", ""))
            == "image_ref"
            for b in blocks
        )
        if not has_inline_refs and img_blocks:
            blocks.extend(img_blocks)
            flat_text = blocks_to_text(blocks)

        return ExtractionResult(
            doc_id=task["doc_id"],
            version_no=task["version_no"],
            source_key=task.get("raw_key", ""),
            file_ext="docx",
            extract_method="python_docx",
            title=title,
            text=flat_text,
            text_length=len(flat_text),
            blocks=blocks,
            warnings=warnings,
            assets=assets,
            vlm_degraded_count=vlm_degraded,
            # E：判决记录不可用时的降级留痕。⚠️ DOCX 分支此前**完全没有** partial_loss_notes
            # 通道 —— 同一处 A11 的 unreadable 留痕返回值也被丢掉（既有缺口，未在本次修，
            # 补它会让"含坏图的 DOCX"新落 NEEDS_REVIEW，是独立的行为变更）。
            partial_loss_notes=_verdict_notes,
        )

    # ── XLSX / XLS ──

    # 设备清扫基准书：通用部位关键词 fallback 白名单
    _PART_KEYWORDS_FALLBACK = {
        "进片口", "外侧", "电刷", "轴轮", "链条", "齿轮", "丝杆", "油位",
        "温度", "剥离条", "电机", "烘箱", "油壸", "管路", "底盘", "配电箱",
        "变频器", "液压站", "三辊", "涂油辊", "硅油", "粉碎机", "油缸",
        "主机", "仪表", "标识", "外观", "送片组件", "拉伸总成", "出杯口",
        "机顶盖", "活动件", "油路",
    }

    # 表头关键词：命中 ≥2 个则认为是表头行
    _HEADER_KEYWORDS = {"清扫部位名称", "点检部位", "部位名称", "清扫基准",
                        "清扫方法", "清扫工具", "清扫周期", "点检项目",
                        "点检方法", "判定标准", "序号", "类别",
                        "运转中", "停机时", "安全注意事项", "责任人",
                        "所需时间", "频次", "异常处理", "点检人"}

    # 部位名称列：优先匹配这些列名
    _PART_COL_NAMES = {"清扫部位名称", "点检部位", "部位名称"}

    def _extract_xlsx(self, task: dict) -> ExtractionResult:
        """Excel 提取（文本 + 嵌入图片）：逐 sheet 逐行读取单元格文本。

        增强功能（设备清扫基准书类文档）：
        - 子 section 检测："清扫时要点检的项目" 插入 heading 分隔清扫区和点检区
        - 表头行识别 + 部位名称列自动提取 part_candidates
        - 表头行和设备信息行标记 row_role="metadata"，不进入 row card chunk
        - 图片 part_labels 提取（优先匹配 part_candidates，fallback 白名单）
        """
        import re
        from opensearch_pipeline.extraction.image_extraction_utils import extract_images_from_xlsx

        local_path = task.get("local_path", "")
        file_ext = task.get("file_ext", "xlsx").lower()
        blocks = []
        warnings = []
        partial_notes = []  # 批次6：中途异常留痕（有产出但不完整 → NEEDS_REVIEW 收尾）
        all_part_candidates: set = set()  # 跨 sheet 收集所有部位名称

        _img_wb = None  # perf#62：普通模式 workbook 供图片提取复用（read_only 无 _images，不可共享）
        try:
            import openpyxl
            # read_only=False 才能读到合并单元格几何（merged_cells），而合并单元格常承载
            # 分组键（清扫基准书的"类别"列、规格书的分区）。
            # 兜底按"行数"而非文件大小判断：图文型 xlsx（清扫基准书/规格书）文件可达数十 MB，
            # 但只有几十行，应当展开合并；真正的数据导出表（数十万~百万行）才保持 read_only 防止内存暴涨。
            # perf#62：probe 与正式 load 合一——判定走 read_only 时直接复用 probe workbook
            # （read_only 的 iter_rows 每次调用都从 archive 重新流式解析，探过 max_row 不影响
            # 后续正文遍历）；判定走普通模式才二次 load，避免同一文件被完整解析两次。
            wb = None
            _use_ro = True
            try:
                _fsz = os.path.getsize(local_path)
                _probe = openpyxl.load_workbook(local_path, read_only=True, data_only=True)
                _maxr = max((ws.max_row or 0) for ws in _probe.worksheets) if _probe.worksheets else 0
                _use_ro = (_maxr > 50000) or (_fsz > 100 * 1024 * 1024)
                if _use_ro:
                    wb = _probe   # 大表：probe 即正式 workbook，省一次完整 load
                else:
                    _probe.close()
            except Exception:
                _use_ro = True
            if wb is None:
                wb = openpyxl.load_workbook(local_path, read_only=_use_ro, data_only=True)
            if not _use_ro:
                # 普通模式加载的 _images 数据在 load 时已读入内存（BytesIO），与 archive
                # 开闭无关，后续 wb.close()（普通模式为 no-op）不影响图片提取复用。
                _img_wb = wb
            for sheet_idx, sheet_name in enumerate(wb.sheetnames):
                ws = wb[sheet_name]

                # 隐藏 sheet 入库与否走证据制（B2-1 条件化，release-gate 2026-07-06 复盘）：
                # 带图/大二维表的隐藏正文 sheet 收录，草稿/下拉源仍跳过；veryHidden 恒跳过。
                # 规则与三档开关 RAG_XLSX_HIDDEN_SHEET_MODE 见 hidden_sheet_should_ingest
                # docstring。不静默——跳过/收录隐藏 sheet 都逐 sheet loud log。enumerate
                # 不跳号 → 可见 sheet 的 sheet_idx/page_num 与图片路径(共用同一决策)保持一致。
                from opensearch_pipeline.extraction.image_extraction_utils import (
                    hidden_sheet_should_ingest,
                )
                _state = getattr(ws, "sheet_state", "visible")
                if not hidden_sheet_should_ingest(ws):
                    warnings.append(f"skipped {_state} sheet: {sheet_name}")
                    print(f"      ⏭️ [xlsx] skip {_state} sheet '{sheet_name}' (not ingested)", flush=True)
                    continue
                if _state != "visible":
                    warnings.append(f"ingested {_state} sheet (body evidence): {sheet_name}")
                    print(f"      👁️ [xlsx] ingest {_state} sheet '{sheet_name}' (body evidence: "
                          f"images/2-D content mass)", flush=True)

                # 合并单元格"向下填充"：纵向合并的首格值（如"类别"列）传播到该列下方各行，
                # 避免下游丢失分组键。横向合并不填充（标题只出现一次）。
                merge_fill = {}
                if not _use_ro:
                    try:
                        for mr in ws.merged_cells.ranges:
                            tl = ws.cell(mr.min_row, mr.min_col).value
                            if tl is None:
                                continue
                            # 只向下传播"短标签"型合并值（分组键，如类别名"设备本体"）。
                            # 跨行合并的长正文/步骤说明若被复制，会生成重复数据行，
                            # 破坏 step 检测与图文绑定（如过程检验 SOP 的步骤被复制两次）。
                            if len(str(tl).strip()) > 20:
                                continue
                            for rr in range(mr.min_row + 1, mr.max_row + 1):
                                merge_fill[(rr, mr.min_col)] = tl
                    except Exception:
                        pass

                # sheet 标题作为 heading block（默认 section_type=cleaning_items）
                sheet_heading = ExtractedBlock(
                    block_type="heading",
                    text=sheet_name,
                    page_num=sheet_idx + 1,
                    section_path=sheet_name,
                    source="openpyxl",
                )
                sheet_heading.extra = {"section_type": "cleaning_items"}
                blocks.append(sheet_heading)

                # ── Pass 1: 扫描所有行，识别表头 + 收集 part_candidates ──
                all_rows = []
                header_row_idx = None
                part_col_idx = None

                for row_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
                    cells = []
                    for ci0, c in enumerate(row):
                        if c is None and (row_idx, ci0 + 1) in merge_fill:
                            c = merge_fill[(row_idx, ci0 + 1)]
                        cells.append(_xlsx_cell_to_str(c) if c is not None else "")
                    all_rows.append((row_idx, cells))

                    # 检测表头行（命中 ≥2 个关键词）
                    if header_row_idx is None:
                        stripped_cells = {c.strip() for c in cells if c.strip()}
                        hits = len(stripped_cells & self._HEADER_KEYWORDS)
                        if hits >= 2:
                            header_row_idx = row_idx
                            # 找部位名称列
                            for ci, c in enumerate(cells):
                                if c.strip() in self._PART_COL_NAMES:
                                    part_col_idx = ci
                                    break

                # 从数据行收集 part_candidates
                if header_row_idx is not None and part_col_idx is not None:
                    for row_idx, cells in all_rows:
                        if row_idx <= header_row_idx:
                            continue
                        if part_col_idx < len(cells):
                            part_name = cells[part_col_idx].strip()
                            if part_name and len(part_name) >= 2 and not part_name.isdigit():
                                all_part_candidates.add(part_name)

                # ── Pass 2: 生成 blocks ──
                # 签字行/备注分界正则
                _RE_SIGNATURE = re.compile(r"编写[：:]?\s*(.*\s+)?审核")
                _RE_SECTION_BOUNDARY = re.compile(r"清扫时[要需]点检|^点检项目$")

                in_inspection_section = False
                for row_idx, cells in all_rows:
                    line = "\t".join(cells).strip()
                    if not line:
                        continue

                    clean_text = line.replace("\t", "").strip()
                    stripped_cells = {c.strip() for c in cells if c.strip()}

                    # ── Bug fix 1: 子区域分界行 → 只生成 heading，跳过 paragraph ──
                    if _RE_SECTION_BOUNDARY.match(clean_text):
                        if not in_inspection_section:
                            in_inspection_section = True
                            sub_heading = ExtractedBlock(
                                block_type="heading",
                                text=f"{sheet_name} — 点检项目",
                                page_num=sheet_idx + 1,
                                section_path=f"{sheet_name} — 点检项目",
                                source="openpyxl",
                            )
                            sub_heading.extra = {"section_type": "inspection_items"}
                            blocks.append(sub_heading)
                        continue  # 不生成 paragraph block

                    # ── Bug fix 3: 签字行 → 跳过 ──
                    if _RE_SIGNATURE.search(clean_text):
                        continue

                    # ── Bug fix 2: 二级表头检测（点检区表头等） ──
                    header_hits = len(stripped_cells & self._HEADER_KEYWORDS)
                    is_secondary_header = (
                        header_row_idx is not None
                        and row_idx > header_row_idx
                        and header_hits >= 2
                    )

                    # ── Bug fix 4: 稀疏表头碎片行（只含表头关键词，无实质数据） ──
                    non_empty_cells = [c.strip() for c in cells if c.strip()]
                    is_sparse_header_fragment = (
                        len(non_empty_cells) <= 3
                        and header_hits >= 1
                        and all(c in self._HEADER_KEYWORDS or len(c) <= 2 for c in non_empty_cells)
                    )

                    # ── v2 fix 2: 纯序号空行（只有 1-2 个 cell 且全是数字/空） ──
                    is_empty_number_row = (
                        len(non_empty_cells) <= 2
                        and all(c.isdigit() for c in non_empty_cells)
                    )

                    # 生成 paragraph block
                    blk = ExtractedBlock(
                        block_type="paragraph",
                        text=line,
                        page_num=sheet_idx + 1,
                        source="openpyxl",
                    )
                    extra = {"row_num": row_idx, "sheet_idx": sheet_idx}

                    # 标记 row_role
                    if header_row_idx is not None and row_idx <= header_row_idx:
                        extra["row_role"] = "metadata"
                    elif is_secondary_header or is_sparse_header_fragment:
                        extra["row_role"] = "metadata"
                    elif is_empty_number_row:
                        extra["row_role"] = "metadata"  # 纯序号空行
                    else:
                        extra["row_role"] = "data"
                        # 提取 part_name（如果有部位列）
                        if part_col_idx is not None and part_col_idx < len(cells):
                            pn = cells[part_col_idx].strip()
                            if pn and len(pn) >= 2 and not pn.isdigit():
                                extra["part_name"] = pn
                        # v2 fix 5: 标记稀疏数据行（非空数据 cell ≤3）
                        data_cells = [c for c in non_empty_cells if not c.isdigit()]
                        if len(data_cells) <= 2 and not is_empty_number_row:
                            extra["sparse_row"] = True

                    blk.extra = extra
                    blocks.append(blk)

            wb.close()
        except Exception as e:
            warnings.append(f"Failed to extract Excel file: {e}")
            # 批次6（ultra unified_extractor:1343）：中途异常保留了已抽取 sheet 的 blocks——
            # 有产出时 0-chunk 疑似失败守卫不查 warnings，文档会以"完整"定稿 DONE、余下
            # sheet 永久缺失。留痕 → 收尾 NEEDS_REVIEW（可服务、重灌自愈）。零产出时
            # 0-chunk 守卫照旧接管（FAILED+retry），本留痕不参与该分支。
            partial_notes.append(f"xlsx_partial: extraction aborted mid-workbook ({e}); "
                                 f"sheets after the failure point are missing")

        # 提取嵌入图片 → 三阶段过滤漏斗
        # perf#62：普通模式 wb 直接传给 extract_images_from_xlsx 复用，省第三次 zip 解包；
        # read_only / 加载失败时 _img_wb=None → 其内部照旧自行加载（graceful degradation）。
        _funnel_stats: dict = {}
        assets, img_blocks, vlm_degraded = self._process_embedded_images(
            extract_images_from_xlsx(local_path, _safe_image_output_dir(task),
                                     workbook=_img_wb),
            task, stats=_funnel_stats,
        )
        partial_notes.extend(_warn_unreadable_images(warnings, _funnel_stats))   # A11（XLSX）
        partial_notes.extend(_warn_verdict_store_unavailable(warnings, _funnel_stats))   # E

        # ── part_labels 提取（混合策略：part_candidates 优先 + 白名单 fallback）──
        if all_part_candidates or self._PART_KEYWORDS_FALLBACK:
            for asset in assets:
                ocr = asset.get("ocr_text", "")
                vs = asset.get("visual_summary", "")
                search_text = f"{ocr} {vs}"
                if not search_text.strip():
                    continue
                # 优先匹配表格中实际出现的部位名称
                labels = [p for p in all_part_candidates if p in search_text]
                # Fallback：通用白名单
                if not labels:
                    labels = [kw for kw in self._PART_KEYWORDS_FALLBACK if kw in search_text]
                if labels:
                    asset["part_labels"] = sorted(set(labels))

        if img_blocks:
            blocks.extend(img_blocks)

        # ── procedure_image_guide 后处理：步骤标注 + 图号映射 ──
        from opensearch_pipeline.extraction.xlsx_classifier import classify_xlsx_layout
        _sheet_names = list(dict.fromkeys(
            b.text for b in blocks
            if b.block_type == "heading" and b.extra.get("section_type") in ("cleaning_items", "")
        ))
        _layout_type, _ = classify_xlsx_layout(
            filename=task.get("filename", ""),
            sheet_names=_sheet_names,
            flat_text=blocks_to_text(blocks)[:5000],
        )
        if _layout_type == "procedure_image_guide":
            _RE_STEP_PREFIX = re.compile(r"^(\d+)\t")
            # 图号引用：如图1、见左图3、4、5、如图8-9
            _RE_FIG_REF = re.compile(
                r"(?:如图|见.*?图)\s*(\d+(?:\s*[、,，\-~～]\s*\d+)*)"
            )

            def _parse_figure_refs(text: str) -> list:
                """从文本中提取所有引用的图号列表，如 ['图1','图3','图4','图5']"""
                refs = []
                for m in _RE_FIG_REF.finditer(text):
                    nums_str = m.group(1)
                    # 处理范围（8-9）和列表（3、4、5）
                    parts = re.split(r"[、,，]\s*", nums_str)
                    for part in parts:
                        part = part.strip()
                        range_m = re.match(r"(\d+)\s*[\-~～]\s*(\d+)", part)
                        if range_m:
                            for n in range(int(range_m.group(1)), int(range_m.group(2)) + 1):
                                refs.append(f"图{n}")
                        elif part.isdigit():
                            refs.append(f"图{part}")
                return refs

            # 标记步骤行
            for blk in blocks:
                if blk.block_type != "paragraph":
                    continue
                m = _RE_STEP_PREFIX.match(blk.text)
                if m:
                    step_no = int(m.group(1))
                    blk.extra["step_no"] = step_no
                    blk.extra["row_role"] = "step"
                    fig_refs = _parse_figure_refs(blk.text)
                    if fig_refs:
                        blk.extra["figure_refs"] = fig_refs

            # 建立图号 → asset 映射
            # 关键：跳过 logo/表头装饰图（anchor_row 在图例区之前的图片）
            # 优先用"图例"行号定位图片区起始；其次用 step_row-3；保底 row<5 过滤 logo
            tu_li_rows = [b.extra.get("row_num", 999) for b in blocks
                          if b.block_type == "paragraph" and "图例" in b.text]
            step_rows = [b.extra.get("row_num", 999) for b in blocks if b.extra.get("step_no")]
            if tu_li_rows:
                figure_start_row = min(tu_li_rows)  # "图例" 行即是图片区开始
            elif step_rows:
                figure_start_row = max(0, min(step_rows) - 3)
            else:
                figure_start_row = 5  # 保底：前5行通常是标题/签署区

            sorted_assets = sorted(assets, key=lambda a: a.get("filename", ""))
            figure_map = {}
            fig_counter = 1
            for asset in sorted_assets:
                anchor_row = asset.get("anchor_row")
                # 跳过步骤区之前的图片（logo、表头装饰等）
                if anchor_row is not None and anchor_row < figure_start_row:
                    asset["figure_no"] = None  # 标记为非步骤图片
                    continue
                fig_label = f"图{fig_counter}"
                asset["figure_no"] = fig_label
                figure_map[fig_label] = fig_counter - 1
                fig_counter += 1

        flat_text = blocks_to_text(blocks)
        title = extract_title_from_blocks(blocks, fallback=task.get("filename", ""))

        return ExtractionResult(
            doc_id=task["doc_id"],
            version_no=task["version_no"],
            source_key=task.get("raw_key", ""),
            file_ext=file_ext,
            extract_method="openpyxl",
            title=title,
            text=flat_text,
            text_length=len(flat_text),
            blocks=blocks,
            warnings=warnings,
            assets=assets,
            # 持久化 DAG1 的 layout 判定（用真实 task filename 分类），DAG2 直接消费，不再重分类。
            xlsx_layout_type=_layout_type,
            vlm_degraded_count=vlm_degraded,
            partial_loss_notes=partial_notes,
        )

    # ── PPTX ──

    # 标题占位符 type 值（PP_PLACEHOLDER 枚举的整数值）
    # TITLE=1, CENTER_TITLE=3, VERTICAL_TITLE=5
    _PPTX_TITLE_TYPES = {1, 3, 5}
    # SUBTITLE=4
    _PPTX_SUBTITLE_TYPES = {4}

    def _extract_pptx(self, task: dict) -> ExtractionResult:
        """
        PPTX 提取（Shape 级别结构化 + 表格 + Speaker Notes + 嵌入图片）。

        每个 slide 的 shapes 按类型生成不同 block：
          - Title/Center Title placeholder → heading block (level=1)
          - Subtitle placeholder → heading block (level=2)
          - Body/Content text → paragraph block
          - Table → table block (markdown pipe format)
          - Speaker notes → paragraph block (标记 source="speaker_notes")

        section_path 追踪：slide 标题 → 后续 shapes 继承。
        """
        from opensearch_pipeline.extraction.image_extraction_utils import extract_images_from_pptx

        local_path = task.get("local_path", "")
        blocks = []
        warnings = []
        partial_notes = []  # 批次6：中途异常留痕（有产出但不完整 → NEEDS_REVIEW 收尾）
        slide_count = 0

        try:
            from pptx import Presentation
            prs = Presentation(local_path)
            slide_count = len(prs.slides)

            current_section = None  # section_path 追踪

            for slide_idx, slide in enumerate(prs.slides):
                page_num = slide_idx + 1
                slide_title = None  # 当前 slide 的标题
                # perf#89：记录本 slide 的 blocks 起始下标。page_num 每 slide 唯一、块按序
                # 追加（notes 在 Phase 3、图片块在循环结束后才 extend），故本 slide 的块
                # 恒为 blocks[slide_start:] —— Phase 2 免去每 slide 对全量 blocks 的扫描。
                slide_start = len(blocks)

                # ── Phase 1: 按 shape 类型生成 blocks（G13：递归展开组合形状）──
                for shape in _iter_pptx_shapes(slide.shapes):

                    # ── 表格 shape → table block ──
                    if shape.has_table:
                        rows_text = []
                        for row in shape.table.rows:
                            cells = [cell.text.strip().replace("\n", " ") if cell.text else ""
                                     for cell in row.cells]
                            if any(cells):
                                rows_text.append("| " + " | ".join(cells) + " |")
                        if rows_text:
                            blocks.append(ExtractedBlock(
                                block_type="table",
                                text="\n".join(rows_text),
                                page_num=page_num,
                                section_path=current_section,
                                source="python_pptx",
                                extra={
                                    "table_index": 0,
                                    "row_count": len(rows_text),
                                    "slide_num": page_num,
                                },
                            ))
                        continue  # table shape 已处理，跳过文本提取

                    # ── 文本类 shape ──
                    if not hasattr(shape, "text") or not shape.text.strip():
                        continue

                    text = shape.text.strip()

                    # 判断是否是标题占位符
                    is_title = False
                    is_subtitle = False
                    if shape.is_placeholder:
                        try:
                            ph_type = shape.placeholder_format.type
                            # ph_type 是 PP_PLACEHOLDER 枚举，int(ph_type) 获取整数值
                            ph_type_int = int(ph_type) if ph_type is not None else -1
                            if ph_type_int in self._PPTX_TITLE_TYPES:
                                is_title = True
                            elif ph_type_int in self._PPTX_SUBTITLE_TYPES:
                                is_subtitle = True
                        except Exception:
                            pass

                    if is_title:
                        slide_title = text
                        current_section = text
                        blocks.append(ExtractedBlock(
                            block_type="heading",
                            text=text,
                            level=1,
                            page_num=page_num,
                            section_path=current_section,
                            source="python_pptx",
                            extra={"placeholder": "title", "slide_num": page_num},
                        ))
                    elif is_subtitle:
                        blocks.append(ExtractedBlock(
                            block_type="heading",
                            text=text,
                            level=2,
                            page_num=page_num,
                            section_path=current_section,
                            source="python_pptx",
                            extra={"placeholder": "subtitle", "slide_num": page_num},
                        ))
                    else:
                        # 普通文本框 / Body 占位符 → paragraph block
                        blocks.append(ExtractedBlock(
                            block_type="paragraph",
                            text=text,
                            page_num=page_num,
                            section_path=current_section,
                            source="python_pptx",
                        ))

                # ── Phase 2: 如果没有检测到标题占位符，用 slide 序号做 fallback heading ──
                if slide_title is None:
                    # 检查此 slide 是否有任何内容 block（perf#89：本 slide 首块即
                    # blocks[slide_start]，无新增块时与原全量扫描同为 no-op——
                    # current_section 保持沿用上一 slide 的值）
                    if len(blocks) > slide_start:
                        # 取本 slide 第一个 block 的前 50 字符作为 slide 标题
                        first_text = blocks[slide_start].text[:50]
                        fallback_title = f"Slide {page_num}" + (f": {first_text}" if first_text else "")
                        current_section = fallback_title

                # ── Phase 3: Speaker notes → paragraph block（默认关闭，B2-2）──
                # 演讲备注常含演讲提示/内部备注/旧内容/"此页不要展示"等，默认不入库。需要时
                # 设 RAG_PPTX_INCLUDE_NOTES=1/true/yes 显式开启（沿用项目 default-OFF flag 约定）。
                _include_notes = os.environ.get("RAG_PPTX_INCLUDE_NOTES", "").lower() in ("1", "true", "yes")
                try:
                    if _include_notes and slide.has_notes_slide and slide.notes_slide:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_text and len(notes_text) > 5:
                            blocks.append(ExtractedBlock(
                                block_type="paragraph",
                                text=notes_text,
                                page_num=page_num,
                                section_path=current_section,
                                source="speaker_notes",
                                extra={"slide_num": page_num, "is_notes": True},
                            ))
                except Exception:
                    pass  # notes 提取失败不影响主流程

        except Exception as e:
            warnings.append(f"Failed to extract PPTX text: {e}")
            # 批次6（同 xlsx）：中途异常保留已抽取 slide 的 blocks——有产出时会以"完整"
            # 定稿 DONE、余下 slide 永久缺失。留痕 → NEEDS_REVIEW 收尾（重灌自愈）。
            partial_notes.append(f"pptx_partial: extraction aborted mid-deck ({e}); "
                                 f"slides after the failure point are missing")

        # 提取嵌入图片 → 三阶段过滤漏斗
        img_stats: dict = {}
        _funnel_stats: dict = {}
        assets, img_blocks, vlm_degraded = self._process_embedded_images(
            extract_images_from_pptx(local_path, _safe_image_output_dir(task), stats=img_stats),
            task, stats=_funnel_stats,
        )
        _warn_skipped_vector_images(warnings, img_stats)
        partial_notes.extend(_warn_unreadable_images(warnings, _funnel_stats))   # A11（PPTX）
        partial_notes.extend(_warn_verdict_store_unavailable(warnings, _funnel_stats))   # E
        if img_blocks:
            blocks.extend(img_blocks)

        flat_text = blocks_to_text(blocks)
        title = extract_title_from_blocks(blocks, fallback=task.get("filename", ""))

        return ExtractionResult(
            doc_id=task["doc_id"],
            version_no=task["version_no"],
            source_key=task.get("raw_key", ""),
            file_ext="pptx",
            extract_method="python_pptx",
            title=title,
            text=flat_text,
            text_length=len(flat_text),
            blocks=blocks,
            page_count=slide_count,
            warnings=warnings,
            assets=assets,
            vlm_degraded_count=vlm_degraded,
            partial_loss_notes=partial_notes,
        )

    # ── Plain text / Markdown ──

    def _extract_text(self, task: dict) -> ExtractionResult:
        """纯文本/Markdown/HTML/CSV 提取。"""
        local_path = task.get("local_path", "")
        # 批次9（ultra P3 unified_extractor:1658）：与 dispatcher（extract:819）同一归一
        # （strip+lstrip('.')）——此前只 lower：file_ext=".html"/" csv " 会错过 HTML 去标签
        # /CSV 解析分支，原始标签/逗号流直接进切块。
        file_ext = task.get("file_ext", "txt").lower().strip().lstrip(".")

        try:
            raw_text, warnings = _read_text_with_fallback(local_path)
        except Exception as e:
            return ExtractionResult(
                doc_id=task["doc_id"],
                version_no=task["version_no"],
                source_key=task.get("raw_key", ""),
                file_ext=file_ext,
                extract_method="plain_text",
                title=task.get("filename", ""),
                text="",
                text_length=0,
                warnings=[f"Failed to read file: {e}"],
            )

        # HTML/CSV 此前被直读为纯文本：HTML 满屏标签、CSV 引号/转义糊在一起。
        # 先转成可读文本再走统一分块；转换失败回退原文（保持优雅降级约定）。
        extract_method = "plain_text"
        if file_ext in ("html", "htm"):   # .htm 与 .html 同样去标签（否则 .htm 满屏标签进索引）
            raw_text = _html_to_text(raw_text)
            extract_method = "html_text"
        elif file_ext == "csv":
            raw_text = _csv_to_text(raw_text)
            extract_method = "csv_table"

        blocks = extract_text_file(raw_text, source="native")
        flat_text = blocks_to_text(blocks)
        title = extract_title_from_blocks(blocks, fallback=task.get("filename", ""))

        return ExtractionResult(
            doc_id=task["doc_id"],
            version_no=task["version_no"],
            source_key=task.get("raw_key", ""),
            file_ext=file_ext,
            extract_method=extract_method,
            title=title,
            text=flat_text,
            text_length=len(flat_text),
            blocks=blocks,
            warnings=warnings,
        )

    # ── 嵌入图片通用处理 ──

    # ── 跨文档 VLM 结果持久化缓存（SQLite 化，实现在 opensearch_pipeline/vlm_cache.py）──
    # 本地存储: scratch/vlm_cache.sqlite3（存量 scratch/vlm_cache.json 首开自动迁移、
    #           原样保留供 tests/eval 脚本独立使用）
    # OSS 镜像: processing/cache/vlm_cache.sqlite3（跨 DataWorks 运行 warm-start；
    #           遗留 processing/cache/vlm_cache.json 在本地/镜像两级皆空时兜底迁移）
    # 缓存 key: "{图片SHA-256}:{ns}"（B4 起主键由 MD5 换 SHA-256；遗留裸 MD5 只读兼容
    #           见 _vlm_cache_lookup）
    # 缓存 value: funnel_result dict (status, visual_summary, reason, width, height, ...)
    #            形态由 _vlm_cache_entry_usable 单点把关（D1）：非法条目读侧视同未命中
    # 这里只留三个薄委托：类方法名/调用点（pipeline_nodes finally 的 flush、
    # _process_embedded_images 的 load/save）保持不变。

    @classmethod
    def _load_vlm_cache(cls) -> dict:
        """返回进程级共享的 VLM 缓存 dict 门面（get/[]= 命中语义与旧 JSON dict 一致）。"""
        from opensearch_pipeline.vlm_cache import get_vlm_cache
        return get_vlm_cache()

    @classmethod
    def _save_vlm_cache(cls, force_oss: bool = False):
        """per-doc 保存：脏键增量落 sqlite（WAL）+ OSS 镜像降频推送（perf#27 节奏不变）。"""
        from opensearch_pipeline.vlm_cache import save_vlm_cache
        save_vlm_cache(force_oss=force_oss)

    @classmethod
    def flush_vlm_cache_to_oss(cls):
        """run 结束的 OSS flush（perf#27）：把镜像账目未满一轮的尾巴推上去。
        node_extract_text_with_ocr finally 调用；无账/simulate 全程零动作。"""
        from opensearch_pipeline.vlm_cache import flush_vlm_cache_to_oss
        flush_vlm_cache_to_oss()

    def _fetch_funnel_verdicts(self, hashes, is_public):
        """E：按文档批量预读判决记录 → `(三态, {sha256: row})`。

        自开自关一条只读连接：抽取阶段没有现成的 DB 句柄，而把连接生命周期挂到
        extractor 实例上会跨文档并发共享（`pipeline_nodes.py:374`）。每篇一次、只读、
        短连接，代价可忽略。
        """
        conn = None
        try:
            from opensearch_pipeline.db import _get_db_conn
            conn = _get_db_conn()
            return funnel_verdict_store.fetch_many(
                conn, hashes, is_public, effective_funnel_policy_version())
        except Exception as e:
            print(f"    ⚠️ [verdict-store] 取连接失败（降级为纯缓存行为）: {type(e).__name__}: {e}")
            return funnel_verdict_store.UNAVAILABLE, {}
        finally:
            if conn is not None:
                try:
                    conn.close()
                except Exception:
                    pass

    def _record_funnel_verdicts(self, hash_to_result, hash_to_cached, existing_rows,
                                is_public, doc_id, task):
        """E：把本篇判决写进记录（first-writer-wins，已有行不覆盖）。

        缓存命中的图也提升进来（`provenance='cache_promoted'`）—— 在"不重灌存量"的前提
        下，L2 只能这样从存量缓存渐进建起来。但**绝不拿当前模型冒充历史事实**：提升行的
        `vlm_model` 留 None（存量缓存没有模型/时间/首篇信息）。
        """
        fresh, promoted = [], []
        for h, res in (hash_to_result or {}).items():
            if h in existing_rows:
                continue                       # 已有权威行 ⇒ 不覆盖（first-writer-wins）
            (promoted if h in hash_to_cached else fresh).append(
                {"image_sha256": h, "funnel_res": res})
        if not fresh and not promoted:
            return
        conn = None
        try:
            from opensearch_pipeline.db import _get_db_conn
            conn = _get_db_conn()
            policy = effective_funnel_policy_version()
            vno = task.get("version_no")
            if fresh:
                funnel_verdict_store.record_many(
                    conn, fresh, is_public, policy,
                    vlm_model=getattr(getattr(self.config, "ocr", None), "vlm_model", None),
                    doc_id=doc_id, version_no=vno, provenance="fresh_decision")
            if promoted:
                funnel_verdict_store.record_many(
                    conn, promoted, is_public, policy, vlm_model=None,
                    doc_id=doc_id, version_no=vno, provenance="cache_promoted")
        except Exception as e:
            print(f"    ⚠️ [verdict-store] 判决落库跳过（不影响摄取）: {type(e).__name__}: {e}")
        finally:
            if conn is not None:
                try:
                    conn.close()
                except Exception:
                    pass

    def _process_embedded_images(self, image_assets: list, task: dict, stats: dict = None) -> tuple:
        """
        将提取出的嵌入图片送入 ImageFunnelProcessor 三阶段过滤漏斗（并发 + 去重 + 持久缓存）。

        优化策略：
          1. Funnel 1（静态启发式，<1ms/张）串行预过滤，快速丢弃装饰图
          2. MD5 Hash 去重 + 跨文档持久缓存：
             - 文档内：相同 hash 图片只过一次 VLM
             - 跨文档：命中 scratch/vlm_cache.json 的图片直接复用，无需 VLM 调用
          3. 未命中缓存的唯一图片并发送入 Funnel 2+3（OCR + VLM），
             使用 ThreadPoolExecutor，并发度由 RAG_VLM_CONCURRENCY 控制（默认 8）

        路由结果：
          - DISCARD_DECORATIVE → 丢弃，不记录
          - ROUTE_TO_TEXT → 记录 asset + 追加 OCR 文本块到 blocks
          - ROUTE_TO_VECTOR → 记录 asset（downstream 自动创建 image chunk）
          - QUARANTINE_SENSITIVE → 记录 asset + warning

        Args:
            image_assets: ImageAsset 列表（来自 image_extraction_utils）。
            task: 当前文档的 task dict。

        Returns:
            (assets, ocr_blocks, vlm_degraded_count)：assets 列表、ROUTE_TO_TEXT 产生的
            文本块列表、以及本文档拿到 VLM 降级兜底结论的图片数（P2-32：含 degraded 标记的
            asset 张数 + 漏斗执行异常整张丢失的唯一图片数）。>0 时调用方经 ExtractionResult
            → canonical → stage-2 收尾把文档落 NEEDS_REVIEW（不 DONE），供应商恢复后可定位重扫。
        """
        if not image_assets:
            return [], [], 0

        import hashlib
        import time
        from concurrent.futures import ThreadPoolExecutor, as_completed
        from opensearch_pipeline.image_funnel_processor import ImageFunnelProcessor

        processor = ImageFunnelProcessor(simulate=self.simulate)
        is_public = "_quarantine/" not in task.get("raw_key", "")
        doc_id = task["doc_id"]

        # ── Phase 1: Funnel 1 串行预过滤（<1ms/张，无需并发） ──
        candidates = []  # 通过 Funnel 1 的图片
        discard_count = 0
        unreadable_names = []   # A11：解码失败的图（≠装饰图），经 stats 出参回传给调用方留痕

        for img_asset in image_assets:
            try:
                # A11（2026-07-25）：改调 processor.screen_static —— 此前这里内联了一份与
                # process_image 完全相同的判据（阈值 50/3.0/8.0 各写一遍），两处可以各自漂移；
                # 且因为异常在 _static_heuristics 内部已被吞掉，下面那个 except 其实永远看不到
                # 解码失败，坏图只会被计成 decorative。
                screen = processor.screen_static(img_asset.local_path)
                fname = os.path.basename(img_asset.local_path)
                if screen["verdict"] == "UNREADABLE":
                    print(f"    [Funnel 1] UNREADABLE image: {fname} ({screen['error']})")
                    unreadable_names.append(f"{img_asset.original_name or fname}: {screen['error']}")
                    continue
                if screen["verdict"] == "DECORATIVE":
                    print(f"    [Funnel 1] Discarded decorative image: {fname} "
                          f"({screen['width']}x{screen['height']}, {screen['file_size_kb']:.1f}KB, "
                          f"ratio={screen['aspect_ratio']:.1f})")
                    discard_count += 1
                    continue
                candidates.append(img_asset)
            except Exception as e:
                print(f"      ⚠️ Funnel 1 heuristic failed for {img_asset.original_name}: {e}")
                unreadable_names.append(f"{img_asset.original_name}: screen failed: {e}")
                continue

        if discard_count:
            print(f"      [Funnel 1] Pre-filtered: {discard_count} decorative, {len(candidates)} remaining")
        if unreadable_names:
            print(f"      [Funnel 1] ⚠️ {len(unreadable_names)} image(s) unreadable (decode failed)")
        if stats is not None:
            stats["unreadable_count"] = len(unreadable_names)
            stats["unreadable_names"] = list(unreadable_names)

        if not candidates:
            return [], [], 0

        # ── Phase 1.5: 内容 Hash 去重 + 跨文档缓存查询 ──
        # B4（P2-03）：MD5→SHA-256——此 hash 兼任 VLM **安全判定**缓存主键（CLEAN/
        # QUARANTINE 结论按键复用），MD5 构造碰撞可让恶意图继承他图的 CLEAN 结论；
        # 配合 _vlm_cache_ns 默认版本 "2"，存量 md5 键整体失效不被误命中。
        vlm_cache = self._load_vlm_cache()

        hash_to_candidates = {}   # sha256 -> [img_asset, ...]
        hash_to_representative = {}  # sha256 -> 第一张图片（代表）
        hash_to_cached_result = {}   # sha256 -> funnel_res（来自持久缓存）

        for img_asset in candidates:
            try:
                with open(img_asset.local_path, "rb") as f:
                    file_hash = hashlib.sha256(f.read()).hexdigest()
            except Exception:
                file_hash = f"fallback_{id(img_asset)}"

            if file_hash not in hash_to_candidates:
                hash_to_candidates[file_hash] = []
                hash_to_representative[file_hash] = img_asset
                # 查询跨文档持久缓存（按 is_public 分命名空间）：public 文档会 bypass 安全审计，
                # 其 CLEAN 结论绝不能被 quarantine 文档复用（否则跳过敏感审计）；反之亦然。
                # 遗留裸 MD5 key 的回退/迁移逻辑见 _vlm_cache_lookup。
                cached = _vlm_cache_lookup(vlm_cache, file_hash, is_public)
                if cached is not None:
                    # 反幻觉清洗同样覆盖历史缓存里的编造 ocr_text（幂等；
                    # 缓存条目自带原图 width/height，密度上界可触发）
                    if cached.get("ocr_text"):
                        clean, _m = sanitize_ocr_text(
                            cached.get("ocr_text", ""),
                            width=cached.get("width") or None,
                            height=cached.get("height") or None,
                        )
                        if clean != cached.get("ocr_text"):
                            cached = dict(cached)
                            cached["ocr_text"] = clean
                            vlm_cache[f"{file_hash}:{_vlm_cache_ns(is_public)}"] = cached
                    hash_to_cached_result[file_hash] = cached
            hash_to_candidates[file_hash].append(img_asset)

        total_unique = len(hash_to_representative)
        dup_count = len(candidates) - total_unique
        cache_hit_count = len(hash_to_cached_result)
        need_vlm_hashes = [h for h in hash_to_representative if h not in hash_to_cached_result]

        if dup_count > 0 or cache_hit_count > 0:
            print(f"      [Hash Dedup] {len(candidates)} candidates → {total_unique} unique, "
                  f"{cache_hit_count} cache hits, {len(need_vlm_hashes)} need VLM")

        # ── 单文档图片硬上限（默认无限）──
        # 一个图爆多的文档 = N 张图 × (1 OCR + 1 VLM) DashScope 调用，无任何硬界。RAG_FUNNEL_MAX_IMAGES
        # 给一个可预测的 per-doc 上限：>0 且 need_vlm 超限时，截断尾部并 loud log（绝不静默丢）。这是
        # 与成本熔断器(opt-in、按 RMB)互补的、默认安全的最坏情况界；默认 0 = 不限 → 零行为变化。
        try:
            _funnel_cap = int(os.environ.get("RAG_FUNNEL_MAX_IMAGES", "0"))
        except ValueError:
            _funnel_cap = 0
        # 批次6（ultra unified_extractor:1945）：cap/DENY 跳过的唯一图片计入文档级降级数——
        # 被跳过的 hash 不进 hash_to_result，Phase 3 直接丢（无 asset、无 OCR 块），此前又不
        # 计入 vlm_degraded_count → 文档定稿 DONE，调高预算后无任何重扫信号，图片永久缺失。
        # 计入后文档收尾 NEEDS_REVIEW（与 degraded/漏斗异常同一自愈通道）。
        _budget_skipped_count = 0
        if _funnel_cap > 0 and len(need_vlm_hashes) > _funnel_cap:
            _budget_skipped_count += len(need_vlm_hashes) - _funnel_cap
            print(f"    🚨 [img-funnel] {doc_id}: {len(need_vlm_hashes)} unique images exceed "
                  f"RAG_FUNNEL_MAX_IMAGES={_funnel_cap} → capping; "
                  f"{len(need_vlm_hashes) - _funnel_cap} image(s) skipped (no VLM/OCR)", flush=True)
            need_vlm_hashes = need_vlm_hashes[:_funnel_cap]

        # ── Cost ceiling: 把 always-on 图片漏斗纳入与 VLM-rebuilder 相同的成本熔断器 ──
        # 此前熔断器只 gate 默认关闭的 rebuilder（unified_extractor:434/443），always-on 的嵌入图
        # 漏斗（每张图 = 1 OCR + 1 VLM 调用）完全无成本上限。这里用 *同一个* run 级 breaker 计费：
        # breaker 默认 disabled → no-op（零行为变化）；一旦启用，单文档/单次运行的 VLM 图片预算
        # 同样约束漏斗。DENY → 跳过本文档剩余 VLM（图片不过 VLM），loud log + run 级熔断告警，
        # 绝不静默。估算/预留任何异常都放行漏斗（优雅降级，不阻断抽取）。
        _breaker = getattr(self, "cost_breaker", None)
        if _breaker is not None and getattr(_breaker, "enabled", False) and need_vlm_hashes:
            try:
                from opensearch_pipeline.extraction.cost_breaker import estimate_doc_cost
                _file_ext = (task.get("file_ext")
                             or os.path.splitext(task.get("local_path", ""))[1]).lstrip(".").lower()
                _est = estimate_doc_cost(_file_ext, unit_count=len(need_vlm_hashes),
                                         cached_count=0, cfg=self.config)
                _allowed, _reason = _breaker.try_reserve(doc_id, _est)
                if not _allowed:
                    print(f"    🚨 [CostBreaker] funnel DENY {doc_id}: {_reason} "
                          f"→ skipping VLM for {len(need_vlm_hashes)} image(s) this doc", flush=True)
                    _breaker.maybe_alert_run_tripped()
                    _budget_skipped_count += len(need_vlm_hashes)   # 批次6：DENY 跳过同样计降级
                    need_vlm_hashes = []
            except Exception as _ce:
                print(f"    ⚠️ [CostBreaker] funnel reserve skipped (non-fatal): {_ce}", flush=True)

        # ── Phase 2: Funnel 2+3 并发处理，仅处理未命中缓存的唯一图片 ──
        max_workers = int(os.environ.get("RAG_VLM_CONCURRENCY", "8"))
        assets = []
        ocr_blocks = []
        t0 = time.time()

        def _process_single(img_asset):
            """单张图片的 Funnel 2+3 处理（线程安全）。"""
            return processor.process_image(
                img_asset.local_path, doc_id, is_public=is_public,
                doc_title=_funnel_doc_title(task),
            )

        # 合并结果：缓存命中 + 新处理
        hash_to_result = dict(hash_to_cached_result)  # 先放入缓存命中的

        # ── E：判决记录预读（flag RAG_FUNNEL_VERDICT_STORE 默认 OFF）──
        # 一次 IN(...) 拿全篇，不逐图往返。三态：拿到记录 / 确认没有 / **拿不到**——
        # 第三态不能当成第二态，否则"库一抖就允许重判"，而重判正是要挡的事。
        _vs_state, _vs_rows, _vs_undecided = None, {}, 0
        _vs_declined = []           # pin 被红线拒绝的图（P0-1/P0-2）
        if funnel_verdict_store.store_enabled():
            _vs_state, _vs_rows = self._fetch_funnel_verdicts(
                list(hash_to_representative), is_public)
            if _vs_state == funnel_verdict_store.UNAVAILABLE:
                # 权威判决拿不到 ⇒ 只有同策略缓存能兜底；两边都没有的图，本次是"现场重判"，
                # 必须显式降级留痕（partial_loss_notes → NEEDS_REVIEW），不得静默定稿。
                _vs_undecided = len([h for h in hash_to_representative
                                     if h not in hash_to_cached_result])
                if _vs_undecided and stats is not None:
                    # 走 per-call 的 stats 字典而**不是**实例属性 —— 一个 UnifiedExtractor
                    # 会被跨文档并发复用（pipeline_nodes.py:374），实例状态会串篇。
                    stats["verdict_store_note"] = funnel_verdict_store.unavailable_note(
                        doc_id, _vs_undecided)

        # ── E / 审查 P1-7：**缓存命中的图也要受权威判决约束** ──
        # 此前 pin 只作用于 `need_vlm_hashes` 那条 as_completed 循环 ⇒ 缓存与 059 冲突时
        # **缓存获胜**：某次 RDS 抖动（UNAVAILABLE ⇒ 不 pin）下 VLM 漂移判出的 DISCARD 会被
        # 写进共享缓存；下一轮 RDS 恢复、059 里权威行仍是 keep，但缓存命中直接短路，
        # 权威行从头到尾没被查询 —— 图被丢弃且**零信号**（`_vs_undecided` 只统计不在缓存
        # 里的图，连 NEEDS_REVIEW 都不会有）。多机场景（笔记本与 DataWorks 各持一份本地
        # vlm_cache、共用同一张 059）无需任何故障即可复现。
        # 写侧本来就把缓存命中的结论以 cache_promoted 建记录 —— 读侧不信记录、写侧却拿
        # 缓存去建记录，两侧方向相反。这里补齐读侧。
        if _vs_rows:
            for _h, _cached in list(hash_to_result.items()):
                _rec = _vs_rows.get(_h)
                if _rec and funnel_verdict_store.recordable(_cached):
                    _pinned = funnel_verdict_store.apply_verdict(_cached, _rec)
                    if _pinned.get("verdict_pin_declined"):
                        _vs_declined.append(_h)
                    elif _pinned.get("status") != _cached.get("status"):
                        print(f"      [verdict-store] 缓存条目与权威记录冲突，按记录纠正: "
                              f"{str(_h)[:12]}… {_cached.get('status')} → {_pinned['status']}")
                    hash_to_result[_h] = _pinned

        # P2-32：漏斗执行异常（OCR/VLM 供应商故障导致 process_image 整体 raise）→ 该唯一图片
        # 完全没有结论、不进 assets——与 degraded 兜底同属「供应商降级」，计入文档级降级数，
        # 使文档收尾走 NEEDS_REVIEW、供应商恢复后重扫可自愈（否则静默 INDEXED 终态无物重扫）。
        funnel_exception_count = 0

        if need_vlm_hashes:
            with ThreadPoolExecutor(max_workers=max_workers) as pool:
                future_to_hash = {}
                for file_hash in need_vlm_hashes:
                    representative = hash_to_representative[file_hash]
                    future = pool.submit(_process_single, representative)
                    future_to_hash[future] = file_hash

                for future in as_completed(future_to_hash):
                    file_hash = future_to_hash[future]
                    try:
                        funnel_res = future.result()
                        # E：有权威记录就把 status 盖回去 —— caption 本表刻意不存、只能
                        # 重掷，但**图集不能因此改变**。这正是 E 的主目标。
                        _rec = _vs_rows.get(file_hash)
                        if _rec and funnel_verdict_store.recordable(funnel_res):
                            funnel_res = funnel_verdict_store.apply_verdict(funnel_res, _rec)
                            if funnel_res.get("verdict_pin_declined"):
                                # P0-1/P0-2 红线：pin 被拒 ⇒ 该图不受权威记录保护。
                                # 不能静默 —— 走与 UNAVAILABLE 同一条降级通道。
                                _vs_declined.append(file_hash)
                        hash_to_result[file_hash] = funnel_res
                        # 写入持久缓存：跳过降级结果（VLM 超时/解析失败的兜底，缓存会把一次性
                        # 故障变成跨文档/跨运行的永久错误标签）和无法稳定取哈希的 fallback key
                        # （基于内存地址，跨进程无意义）。缓存键按 is_public 分命名空间。
                        # ⚠️ simulate 模式的 mock 结果绝不能入持久缓存 —— 按真实 MD5 落键后，
                        # 真实运行会命中 mock 描述（缓存投毒）。
                        # A11：DISCARD_UNREADABLE 同样**不得入缓存** —— 与上面 degraded 的理由
                        # 逐字相同：解码失败往往是一次性的（临时文件被截断/落盘竞态/内存压力），
                        # 缓存会把它固化成该图跨文档、跨运行的永久标签，而缓存命中侧不校验状态。
                        if _vlm_cache_put_allowed(self.simulate, funnel_res, file_hash):
                            cache_key = f"{file_hash}:{_vlm_cache_ns(is_public)}"
                            vlm_cache[cache_key] = _vlm_cache_entry(funnel_res)
                    except Exception as e:
                        rep = hash_to_representative[file_hash]
                        funnel_exception_count += 1
                        print(f"      ⚠️ Image funnel failed for {rep.original_name}: {e}")

            # 处理完本文档后持久化缓存（simulate 不落盘，防 mock 污染共享缓存）
            if not self.simulate:
                self._save_vlm_cache()

        # pin 被拒的图同样要留痕（与 UNAVAILABLE 同通道）。**追加不覆盖** —— 两种降级
        # 可以同时发生（RDS 抖动 + 本轮判敏感）。
        if _vs_declined and stats is not None:
            _pn = funnel_verdict_store.pin_declined_note(doc_id, len(_vs_declined))
            stats["verdict_store_note"] = (
                (stats.get("verdict_store_note", "") + " || " + _pn).lstrip(" |")
                if stats.get("verdict_store_note") else _pn)

        # ── E：把本篇的判决落进记录（first-writer-wins，已有行不覆盖）──
        # 缓存命中的图也一并提升（provenance=cache_promoted）：不重灌存量的前提下，L2 只能
        # 这样渐进建起来。但**绝不拿当前模型冒充历史事实** —— 提升行的 vlm_model 留空。
        if _vs_state in (funnel_verdict_store.HIT, funnel_verdict_store.MISS) and not self.simulate:
            self._record_funnel_verdicts(hash_to_result, hash_to_cached_result,
                                         _vs_rows, is_public, doc_id, task)

        # ── Phase 3: 扇出结果到所有图片（包括重复项） ──
        all_results = []  # (img_asset, funnel_res)
        for file_hash, img_assets_group in hash_to_candidates.items():
            if file_hash not in hash_to_result:
                continue
            funnel_res = hash_to_result[file_hash]
            for img_asset in img_assets_group:
                all_results.append((img_asset, funnel_res))

        # 按 page_num → anchor_row → image_index 排序，保持文档内图片的原始顺序
        # DOCX 的 page_num 全部是 None，需要用 image_index 作为 fallback。
        # XLSX 在同一 sheet 内允许 image_index 跨 sheet 累加而 anchor_row 才是物理行号 —
        # 把 anchor_row 提到 image_index 之前，确保 xlsx assets 的下游 pool 顺序由
        # 物理行号唯一决定，避免任何 ThreadPoolExecutor / 历史索引顺序的残留影响。
        all_results.sort(key=lambda x: (
            x[0].page_num if x[0].page_num is not None else 999999,
            getattr(x[0], "anchor_row", None) if getattr(x[0], "anchor_row", None) is not None else 999999,
            x[0].image_index if x[0].image_index is not None else 999999,
            x[0].original_name or "",
        ))

        degraded_asset_count = 0
        # A11：Phase-3 兜底 —— process_image 内部也可能判 UNREADABLE（Phase-1 与它各跑一次
        # 解码，第二次才失败的竞态确实存在），这些同样要计入留痕。
        _unreadable_late = []
        for img_asset, funnel_res in all_results:
            status = funnel_res["status"]

            # Funnel 1 结果在 process_image 内部也会触发（双重保护），跳过
            # A11：DISCARD_UNREADABLE 一并跳过——**必须是精确等值枚举而非 DISCARD 前缀匹配**：
            # 这条链路上确实只有它是等值判断，漏掉就会让读不出来的图混进 asset_dict，
            # 带着 0×0 尺寸和空 ocr_text 一路流进 image_refs 契约。
            if status in ("DISCARD_DECORATIVE", "DISCARD_UNREADABLE"):
                if status == "DISCARD_UNREADABLE":
                    _unreadable_late.append(
                        f"{img_asset.original_name}: {funnel_res.get('reason', 'unreadable')}")
                continue

            asset_dict = {
                "filename": os.path.basename(img_asset.local_path),
                "local_path": img_asset.local_path,
                "page_num": img_asset.page_num,
                "image_index": img_asset.image_index,
                "original_index": img_asset.image_index,
                "status": status,
                "width": funnel_res.get("width", 0),
                "height": funnel_res.get("height", 0),
                "file_size_kb": funnel_res.get("file_size_kb", 0.0),
                "ocr_text": funnel_res.get("ocr_text", ""),
                "visual_summary": funnel_res.get("visual_summary", ""),
                "image_category": funnel_res.get("image_category", ""),

                "vlm_annotation_map": funnel_res.get("vlm_annotation_map", {}),
                "reason": funnel_res.get("reason", ""),
            }
            # 透传 PDF 显示 bbox（页坐标、上原点 — 版面位置图片锚定用）
            if getattr(img_asset, "bbox", None):
                asset_dict["bbox"] = list(img_asset.bbox)
            # 透传 XLSX anchor_row（用于行级图片绑定）
            if hasattr(img_asset, "anchor_row") and img_asset.anchor_row is not None:
                asset_dict["anchor_row"] = img_asset.anchor_row
            # 透传 XLSX annotation_num（Drawing XML 分组标注编号，如 ①②③）
            if hasattr(img_asset, "annotation_num") and img_asset.annotation_num is not None:
                asset_dict["annotation_num"] = img_asset.annotation_num
            if hasattr(img_asset, "annotation_label") and img_asset.annotation_label:
                asset_dict["annotation_label"] = img_asset.annotation_label
            # P2-32：VLM 降级兜底结论（占位 caption / 保守隔离）逐 asset 透传 + 计入文档级降级数
            # （只在为真时落键，不膨胀 canonical JSON）。
            if funnel_res.get("degraded"):
                asset_dict["degraded"] = True
                degraded_asset_count += 1
            assets.append(asset_dict)

            # ROUTE_TO_TEXT：追加 OCR 文本块（与 _extract_image 行为一致）
            if status == "ROUTE_TO_TEXT" and funnel_res.get("ocr_text"):
                ocr_blocks.append(ExtractedBlock(
                    block_type="ocr_text",
                    text=funnel_res["ocr_text"],
                    page_num=img_asset.page_num,
                    source="ocr",
                    extra={
                        "vlm_annotation_map": funnel_res.get("vlm_annotation_map", {}),
                        "visual_summary": funnel_res.get("visual_summary", ""),
                        "image_category": funnel_res.get("image_category", ""),
                        "source_image": os.path.basename(img_asset.local_path),
                        "local_path": img_asset.local_path,
                    },
                ))

        elapsed = time.time() - t0
        routed_counts = {}
        for a in assets:
            s = a["status"]
            routed_counts[s] = routed_counts.get(s, 0) + 1
        if routed_counts:
            summary = ", ".join(f"{k}={v}" for k, v in routed_counts.items())
            vlm_calls = len(need_vlm_hashes)
            avg_ms = (elapsed / vlm_calls * 1000) if vlm_calls else 0
            print(f"      [img-funnel] {len(image_assets)} extracted → {len(assets)} kept ({summary})")
            print(f"      [img-funnel] ⚡ VLM calls={vlm_calls}, cache_hits={cache_hit_count}, "
                  f"dedup={dup_count}, time={elapsed:.1f}s ({avg_ms:.0f}ms/call, workers={max_workers})")

        # A11：把 Phase-3 兜底抓到的 unreadable 并进 stats（Phase-1 已写过一次，这里累加）
        if _unreadable_late and stats is not None:
            stats["unreadable_names"] = list(stats.get("unreadable_names", [])) + _unreadable_late
            stats["unreadable_count"] = len(stats["unreadable_names"])

        # P2-32：文档级降级数 = degraded asset 张数 + 漏斗异常整张丢失的唯一图片数
        # + 批次6：预算 cap/DENY 跳过的唯一图片数（同一 NEEDS_REVIEW 自愈通道）。
        vlm_degraded_count = degraded_asset_count + funnel_exception_count + _budget_skipped_count
        if vlm_degraded_count:
            print(f"      🚨 [img-funnel] {doc_id}: {vlm_degraded_count} 张图片为 VLM 降级/缺失"
                  f"（degraded assets={degraded_asset_count}, funnel exceptions="
                  f"{funnel_exception_count}, budget skipped={_budget_skipped_count}）"
                  f"→ 文档将收尾 NEEDS_REVIEW，供应商恢复/预算调高后重扫自愈")

        return assets, ocr_blocks, vlm_degraded_count

    # ── Image (direct OCR) ──

    def _extract_image(self, task: dict) -> ExtractionResult:
        """图片直接调用过滤漏斗进行三阶段分析。"""
        from opensearch_pipeline.image_funnel_processor import ImageFunnelProcessor
        
        local_path = task.get("local_path", "")
        is_public = "_quarantine/" not in task.get("raw_key", "")
        
        processor = ImageFunnelProcessor(simulate=self.simulate)
        # A13（2026-07-25）：接入跨文档持久 VLM 缓存 —— 此前独立图片文档直接 process_image，
        # 前后无查/无写，每次首灌/升版/重试都重付 1 OCR + 1 VLM。金额是小头，**真正的理由是
        # 确定性**：visual_summary 每次重掷正是 xlsx 图↔步骤绑定漂移那场战役的根因。
        # 缓存语义与嵌入图逐字一致（sha256 主键 + pub/sec 命名空间 + 版本；写入门与条目形态
        # 都走 _vlm_cache_put_allowed / _vlm_cache_entry 单一来源）。
        file_hash = _image_sha256(local_path)
        vlm_cache = self._load_vlm_cache()
        funnel_res = _vlm_cache_lookup(vlm_cache, file_hash, is_public)
        if funnel_res is not None:
            # 命中也要过反幻觉复洗（与嵌入图路径同款，幂等）
            if funnel_res.get("ocr_text"):
                _clean, _m = sanitize_ocr_text(
                    funnel_res.get("ocr_text", ""),
                    width=funnel_res.get("width") or None,
                    height=funnel_res.get("height") or None,
                )
                if _clean != funnel_res.get("ocr_text"):
                    funnel_res = dict(funnel_res)
                    funnel_res["ocr_text"] = _clean
                    vlm_cache[f"{file_hash}:{_vlm_cache_ns(is_public)}"] = funnel_res
            print(f"      [VLM Cache] hit for standalone image {os.path.basename(local_path)}")
        else:
            funnel_res = processor.process_image(
                local_path, task["doc_id"], is_public=is_public,
                doc_title=_funnel_doc_title(task),
            )
            if _vlm_cache_put_allowed(self.simulate, funnel_res, file_hash):
                vlm_cache[f"{file_hash}:{_vlm_cache_ns(is_public)}"] = _vlm_cache_entry(funnel_res)
                self._save_vlm_cache()

        status = funnel_res["status"]
        assets = [{
            "filename": os.path.basename(local_path),
            "local_path": local_path,
            # 与嵌入图 asset_dict 字段对齐：image_index/page_num/funnel 决策可溯源，
            # 下游 _enrich_existing_image_refs 按 image_index 注入 source_image
            "image_index": 0,
            "original_index": 0,
            "page_num": 1,
            # 独立图片文档：raw/ 对象本身就是这张图，直接作为可签名的 oss_key ——
            # 资产上传环节只上传 ROUTE_TO_VECTOR，构造出的 processing/assets/ 路径
            # 对 ROUTE_TO_TEXT 永不存在（serving 签出 403 死图，对抗评审证实）。
            "oss_key": task.get("raw_key", ""),
            "status": status,
            "width": funnel_res.get("width", 0),
            "height": funnel_res.get("height", 0),
            "file_size_kb": funnel_res.get("file_size_kb", 0.0),
            "ocr_text": funnel_res.get("ocr_text", ""),
            "visual_summary": funnel_res.get("visual_summary", ""),
            "image_category": funnel_res.get("image_category", ""),

            "vlm_annotation_map": funnel_res.get("vlm_annotation_map", {}),
            "reason": funnel_res.get("reason", "")
        }]
        # P2-32：独立图片文档的 VLM 降级兜底同样透传（与嵌入图 asset_dict 对齐）
        if funnel_res.get("degraded"):
            assets[0]["degraded"] = True

        # 根据漏斗决策构建块和全文
        blocks = []
        warnings = []

        if status == "ROUTE_TO_TEXT" and funnel_res.get("ocr_text"):
            blocks.append(ExtractedBlock(
                block_type="ocr_text",
                text=funnel_res["ocr_text"],
                page_num=1,
                source="ocr",
                extra={
                    "vlm_annotation_map": funnel_res.get("vlm_annotation_map", {}),
                    "visual_summary": funnel_res.get("visual_summary", ""),
                    "image_category": funnel_res.get("image_category", ""),
                    "local_path": local_path,
                },
            ))
            # 独立图片文档的图就是文档本体（磨床操作流程.png 等 SOP 海报）：
            # ROUTE_TO_TEXT 只留 OCR 文本会让原图永远无法在回答里渲染。
            # 附 image_ref 块（仅 image_index 占位），chunk 阶段由
            # _enrich_existing_image_refs 注入 source_image/visual_summary。
            blocks.append(ExtractedBlock(
                block_type="image_ref",
                text="",
                page_num=1,
                source="multimodal",
                extra={"image_index": 0},
            ))
        elif status == "QUARANTINE_SENSITIVE":
            warnings.append(f"🚨 Sensitive content detected in non-public image asset: {funnel_res.get('reason')}")

        # A11：独立图片文档 = 那张图就是文档本体，读不出来 ⇒ 整篇没有任何内容。
        # 走既有 0-chunk「疑似失败」闭环：warning 文案命中 node_write_chunk_meta 的关键词表
        # （decode failed 里的 "fail"）→ chunk_status='NEEDS_REVIEW' +
        # content_process_status='FAILED' + retry_count+1，沿用毒文档重试预算，3 次后停在
        # FAILED 等人工（控制台徽章「处理失败」）。**不是**自愈：stage-2 重试重跑的是同一份
        # canonical 的切块，不会重新解码图片，三次结果确定性相同。
        _unreadable_notes = []
        if status == "DISCARD_UNREADABLE":
            _unreadable_notes = _warn_unreadable_images(
                warnings,
                {"unreadable_names": [f"{os.path.basename(local_path)}: "
                                      f"{funnel_res.get('reason', 'decode failed')}"]})
        
        flat_text = blocks_to_text(blocks)

        return ExtractionResult(
            doc_id=task["doc_id"],
            version_no=task["version_no"],
            source_key=task.get("raw_key", ""),
            file_ext=task.get("file_ext", "png"),
            extract_method="image_funnel",
            title=task.get("filename", ""),
            text=flat_text,
            text_length=len(flat_text),
            blocks=blocks,
            ocr_required=(status == "ROUTE_TO_TEXT"),
            ocr_status="DONE" if status == "ROUTE_TO_TEXT" else "NOT_REQUIRED",
            warnings=warnings,
            assets=assets,
            vlm_degraded_count=1 if funnel_res.get("degraded") else 0,
            partial_loss_notes=_unreadable_notes,
        )


    # ── Unsupported ──

    def _unsupported(self, task: dict, file_ext: str) -> ExtractionResult:
        """不支持的文件类型。"""
        return ExtractionResult(
            doc_id=task["doc_id"],
            version_no=task["version_no"],
            source_key=task.get("raw_key", ""),
            file_ext=file_ext,
            extract_method=f"unsupported:{file_ext}",
            title=task.get("filename", ""),
            text="",
            text_length=0,
            warnings=[f"Unsupported file type: {file_ext}"],
        )

    # ── OCR fallback logic ──

    def _pages_needing_ocr(self, result: ExtractionResult) -> List[int]:
        """PDF: 返回原生文本不足的页码（1-based）——扫描页 / 坏字体页 / 空白页。

        旧门槛按整文档 text_length 判断，会让"封面有文字、正文是扫描"的多页 PDF
        一页都不 OCR（#1 语料失败模式）。这里逐页统计原生文本，挑出需要 OCR 的页。
        图片块 / 已有 OCR 块不计入"原生文本"。

        page_count<=0 的两种来源：
          (a) 真的 0 页 PDF（不存在）；
          (b) extraction 失败（pdfplumber 拿不到 + pypdf 也挂）——这才是真情况，
              RD 61D861 就是被旧版静默 return [] 漏过的扫描型 PDF。
        保守策略：尝试用别的 PDF 库重拿 page_count；都拿不到则返回 [1] 占位，
        强制 qwen-vl-ocr 至少看第 1 页（它能自己识别 PDF 真实页数）。
        """
        if result.file_ext != "pdf":
            return []

        page_count = result.page_count or 0
        if page_count <= 0:
            # 保守 fallback：尝试用本地 PDF 库重新拿 page_count（graceful degradation）
            local_path = ""
            try:
                # source_key 不一定是本地路径；这里仅能利用 result 自身已有的信息。
                # 没有 local_path 时只能走 [1] 占位。
                local_path = getattr(result, "_local_path", "") or ""
            except Exception:
                local_path = ""

            recovered = 0
            if local_path:
                # 尝试 pypdf
                try:
                    from pypdf import PdfReader  # type: ignore
                    recovered = len(PdfReader(local_path).pages)
                except Exception:
                    try:
                        from PyPDF2 import PdfReader  # type: ignore
                        recovered = len(PdfReader(local_path).pages)
                    except Exception:
                        recovered = 0
                # 再尝试 pdfplumber
                if recovered <= 0:
                    try:
                        import pdfplumber  # type: ignore
                        with pdfplumber.open(local_path) as pdf:
                            recovered = len(pdf.pages)
                    except Exception:
                        recovered = 0

            warn_msg = (
                f"_pages_needing_ocr: page_count<=0 (extraction likely failed); "
                f"conservative OCR fallback engaged (recovered_page_count={recovered or 'unknown'})"
            )
            if isinstance(getattr(result, "warnings", None), list):
                result.warnings.append(warn_msg)

            if recovered <= 0:
                # 拿不到真实页数 → 至少 OCR 第 1 页占位，让 qwen-vl-ocr 自己看图判断
                return [1]
            # 拿到了真实页数：用真实 page_count 走下面正常逐页路径
            result.page_count = recovered
            page_count = recovered

        per_page: Dict[int, int] = {}
        per_page_text: Dict[int, List[str]] = {}
        for b in result.blocks:
            bt = b.get("block_type", "") if isinstance(b, dict) else getattr(b, "block_type", "")
            if bt in ("image_ref", "ocr_text"):
                continue
            pg = b.get("page_num") if isinstance(b, dict) else getattr(b, "page_num", None)
            if pg is None:
                continue
            txt = (b.get("text", "") if isinstance(b, dict) else getattr(b, "text", "")) or ""
            per_page[pg] = per_page.get(pg, 0) + len(txt.strip())
            per_page_text.setdefault(pg, []).append(txt)
        # 只考察原生抽取覆盖过的页（≤ pdf_native_max_pages）；越界页未被抽取，非"缺文本=需OCR"。
        scan_upto = min(page_count, self.pdf_native_max_pages)
        needy: List[int] = []
        garbled: List[int] = []
        for pg in range(1, scan_upto + 1):
            if per_page.get(pg, 0) < self.PER_PAGE_OCR_THRESHOLD:
                needy.append(pg)
            elif not _native_text_quality_ok("\n".join(per_page_text.get(pg, ()))):
                # G15：字符数够但内容是坏字体乱码（(cid:)/PUA/替换符/非常用字符占比超阈）
                # ——旧门槛下这类页永不重 OCR，垃圾直接进切块与向量。
                needy.append(pg)
                garbled.append(pg)
        if garbled:
            if isinstance(getattr(result, "warnings", None), list):
                result.warnings.append(
                    f"[GARBLED_TEXT] 页 {garbled} 原生文本疑似坏字体乱码，已送 OCR 重取；"
                    f"OCR 成功后这些页的原生块将被替换")
            # 供 _apply_ocr_fallback 在 OCR 成功后剔除这些页的乱码原生块
            # （不暴露成 dataclass field，避免序列化噪声——同 _local_path 先例）
            result._garbled_pages = set(garbled)  # type: ignore[attr-defined]
        return needy

    def _needs_ocr(self, result: ExtractionResult) -> bool:
        """判断是否需要 OCR fallback。"""
        if result.file_ext == "pdf":
            # 逐页判断：任一页原生文本不足即触发（只 OCR 那些页）
            return len(self._pages_needing_ocr(result)) > 0
        if result.file_ext in ("png", "jpg", "jpeg", "webp"):
            return result.text_length < self.OCR_THRESHOLD_CHARS
        return False

    def _apply_ocr_fallback(self, task: dict, result: ExtractionResult,
                            needy: Optional[List[int]] = None) -> ExtractionResult:
        """应用 OCR fallback。PDF 只 OCR 文本不足的页，并按 page_num 合并回正确位置。

        needy: 调用方已算好的 PDF 需 OCR 页列表（perf#65 避免 _pages_needing_ocr 重算，
               该函数在 page_count<=0 时会重开 pypdf/pdfplumber 探测，可达秒级）；
               None = 自行计算（独立调用兼容，行为同旧版）。
        """
        local_path = task.get("local_path", "")

        if result.file_ext == "pdf":
            if needy is None:
                needy = self._pages_needing_ocr(result)
            if not needy:
                return result
            ocr_result = self.ocr_client.ocr_pdf(
                local_path, task["doc_id"], self.oss_client, page_nums=needy,
            )
        else:
            ocr_result = self.ocr_client.ocr_image(
                local_path, task["doc_id"], self.oss_client,
            )

        ocr_blocks = ocr_result.to_blocks()
        if ocr_blocks:
            if result.file_ext == "pdf":
                # G12-OCR：先用原生 Pass-1 页眉/页脚词集裁剪 OCR 文本行（OCR 无坐标，
                # 词集是唯一可迁移的页眉知识；纯扫描件词集为空=no-op，fail-open）
                _hf = getattr(result, "_hf_texts", None) or set()
                if _hf:
                    ocr_blocks, _hf_dropped = _strip_hf_lines_from_ocr(ocr_blocks, _hf)
                    if _hf_dropped:
                        msg = (f"[OCR_HF_CROP] 页级 OCR 文本剔除 {_hf_dropped} 行"
                               f"页眉/页脚（G12 OCR 侧，词集 {len(_hf)} 项）")
                        result.warnings.append(msg)
                        print(f"      {msg}", flush=True)
                # G15：乱码页（质量门标记）在 OCR 成功后剔除其原生文本块——否则
                # (cid:)/PUA 垃圾会与 OCR 文本并存进入切块。仅剔除有 OCR 产出的页
                # （fail-open：OCR 没覆盖到的乱码页保留原生块，宁乱勿丢）。
                garbled = getattr(result, "_garbled_pages", None) or set()
                if garbled:
                    ocr_pages = {getattr(b, "page_num", None) for b in ocr_blocks}
                    drop = garbled & ocr_pages
                    if drop:
                        result.blocks = [
                            b for b in result.blocks
                            if not (getattr(b, "page_num", None) in drop
                                    and getattr(b, "source", "") == "native"
                                    and getattr(b, "block_type", "") != "image_ref")
                        ]
                # G6-③ thin 页原生残文去重：needy 页（原生 <PER_PAGE_OCR_THRESHOLD 字符）
                # 的原生文本块与页级 OCR 全页转写并存 = 同内容双份（channel #3）。页级
                # OCR 是该页的权威整页转写——仅当该页 OCR 产出字符数 ≥ 原生残文字符数
                # 时（OCR 严格更完整，fail-open：OCR 偏短则保留原生，宁重勿丢）剔除
                # 该页原生文本块（image_ref/表格块不动）。
                _ocr_chars: Dict[int, int] = {}
                for _b in ocr_blocks:
                    _pg = getattr(_b, "page_num", None)
                    if _pg is not None:
                        _ocr_chars[_pg] = _ocr_chars.get(_pg, 0) + len(getattr(_b, "text", "") or "")
                _native_chars: Dict[int, int] = {}
                for _b in result.blocks:
                    _pg = getattr(_b, "page_num", None)
                    if (_pg in _ocr_chars and getattr(_b, "source", "") == "native"
                            and getattr(_b, "block_type", "") not in ("image_ref", "table")):
                        _native_chars[_pg] = _native_chars.get(_pg, 0) + len(getattr(_b, "text", "") or "")
                _thin_drop = {p for p, n in _native_chars.items() if _ocr_chars.get(p, 0) >= n}
                if _thin_drop:
                    _before = len(result.blocks)
                    result.blocks = [
                        b for b in result.blocks
                        if not (getattr(b, "page_num", None) in _thin_drop
                                and getattr(b, "source", "") == "native"
                                and getattr(b, "block_type", "") not in ("image_ref", "table"))
                    ]
                    msg = (f"[THIN_PAGE_OCR_DEDUP] 页 {sorted(_thin_drop)} 原生残文块 "
                           f"{_before - len(result.blocks)} 个被页级 OCR 全页转写取代（channel #3）")
                    result.warnings.append(msg)
                    print(f"      {msg}", flush=True)
                # G6 双重 OCR 去重：funnel ROUTE_TO_TEXT 已转写过的整页扫描图，
                # 其文本被本次页级 OCR 覆盖时剔除 funnel 块（详见 _drop_double_ocr_dups）。
                result.blocks, _dup_pages = _drop_double_ocr_dups(result.blocks, ocr_blocks)
                if _dup_pages:
                    msg = (f"[DOUBLE_OCR_DEDUP] 页 {sorted(set(_dup_pages))} 的嵌入图 OCR 文本"
                           f"与页级 OCR 近重复，已剔除 {len(_dup_pages)} 个 funnel 文本块"
                           f"（图片资产保留）")
                    result.warnings.append(msg)
                    print(f"      {msg}", flush=True)
                # OCR 的是空白页（原页几乎无块），按 page_num 稳定排序合并保持文档顺序
                merged = list(result.blocks) + list(ocr_blocks)
                merged.sort(key=lambda b: (getattr(b, "page_num", 0) or 0))
                result.blocks = merged
                result.text = blocks_to_text(result.blocks)
            else:
                result.blocks.extend(ocr_blocks)
                if ocr_result.combined_text:
                    result.text = (result.text + "\n\n" + ocr_result.combined_text
                                   if result.text else ocr_result.combined_text)
            result.text_length = len(result.text)

        result.ocr_required = True
        result.ocr_status = ocr_result.status
        result.extract_method = f"{result.extract_method}+ocr_fallback"

        # 批次6（ultra ocr_client:298）：部分页 OCR 失败 → 留痕驱动 NEEDS_REVIEW 收尾。
        # 失败页 text="" 会被 to_blocks 丢弃、聚合状态仍 DONE——不留痕的话缺页内容永久
        # 静默缺失且无重扫信号；文本照常可服务（graceful degradation），只是不定稿 DONE。
        _failed_pgs = [p.page_num for p in getattr(ocr_result, "pages", []) or []
                       if getattr(p, "status", "") == "FAILED"]
        if _failed_pgs and ocr_result.status != "FAILED":
            _note = (f"ocr_partial: {len(_failed_pgs)}/{len(ocr_result.pages)} page(s) failed "
                     f"(pages {_failed_pgs[:10]}) — content of failed pages missing")
            result.partial_loss_notes.append(_note)
            result.warnings.append(_note)

        # B5（2026-07-25）：被 max_ocr_pages 上限砍掉的页同样要留痕。它们**不产生**
        # OCRPageResult，所以上面那条只看 pages 的判据永远看不见它们 —— 长文档尾部内容
        # 不进索引，而 document_version 照常 DONE+INDEXED，运维无法用一条 SQL 圈出重扫对象。
        # 只留痕：**不动上限本身**（那三个页上限是当前唯一的摄取侧成本闸 —— cost_breaker
        # 因 RebuildConfig.enabled=False 恒 no-op；拆闸之前必须先有 breaker）。
        _capped = int(getattr(ocr_result, "skipped_pages", 0) or 0)
        if _capped > 0:
            _cap_note = (f"ocr_page_cap: {_capped} page(s) never OCR'd — capped by "
                         f"max_ocr_pages; their content is missing from the index")
            result.partial_loss_notes.append(_cap_note)
            result.warnings.append(_cap_note)

        return result
