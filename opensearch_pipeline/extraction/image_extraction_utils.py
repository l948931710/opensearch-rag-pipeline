# -*- coding: utf-8 -*-
"""
image_extraction_utils.py — 文档嵌入图片提取工具

从 DOCX / PDF / XLSX 文件中提取嵌入图片到本地临时目录。
返回 ImageAsset 列表，供 ImageFunnelProcessor 三阶段过滤使用。

设计原则：
  - 所有图片都应被处理，不设人为上限（Funnel 1 的 heuristic 过滤会自然淘汰装饰图）
  - MD5 去重：docx 中同一图片可能被多次引用，只处理一次
  - 异常隔离：图片提取失败不影响文本提取主流程，只记录 warning
"""

import hashlib
import os
from dataclasses import dataclass
from typing import List, Optional


@dataclass
class ImageAsset:
    """从文档中提取的单张嵌入图片。"""
    local_path: str                         # 导出到 tmp_dir 的本地文件路径
    page_num: Optional[int] = None          # 所在页码（PDF 可精确提供，DOCX 为 None）
    image_index: int = 0                    # 在文档中的顺序索引
    original_name: str = ""                 # 在文档包内的原始名称（如 media/image3.jpeg）
    # 页面显示 bbox (x0, y0, x1, y1)，页坐标、上原点（PDF/PPTX 可提供；与文本块
    # extra.y0/y1 同坐标系，供按版面位置锚定图片→步骤）
    bbox: Optional[tuple] = None
    # perf#88：导出时已算过的 blob MD5（xlsx 路径回填；别处构造的旧路径为 None 时
    # _enrich_xlsx_annotations 回退逐文件读哈希）。避免配对阶段 O(pics×assets) 重读+重哈希。
    md5: Optional[str] = None


# ═══════════════════════════════════════════════════════════════
# DOCX — python-docx relationship 遍历
# ═══════════════════════════════════════════════════════════════

def _note_skipped_vector(stats, ext: str) -> None:
    """G4：矢量图（WMF/EMF/SVG）跳过计数——此前三处裸 continue 静默丢弃，损失不可见。

    stats 为调用方传入的 dict（None = 不统计，独立调用兼容）；只累计，不改变跳过行为。
    """
    if stats is None:
        return
    stats["skipped_vector_images"] = stats.get("skipped_vector_images", 0) + 1
    stats.setdefault("skipped_vector_formats", set()).add(ext.lstrip("."))


def extract_images_from_docx(
    local_path: str,
    output_dir: str,
    document=None,
    stats: Optional[dict] = None,
) -> List[ImageAsset]:
    """
    从 DOCX 文件中提取所有嵌入图片。

    策略：遍历 document.part.rels 中 reltype 包含 'image' 的关系，
    读取 target_part.blob 获取图片二进制数据。

    DOCX 不提供原生页码信息，所有图片 page_num 为 None。

    Args:
        local_path: DOCX 文件的本地路径。
        output_dir: 图片导出目标目录。
        document: 可选，调用方已解析好的 python-docx Document（perf#63 复用，
            避免同一 DOCX 再做一次完整 OOXML 解包 + lxml 树构建）；
            None = 自行加载（独立调用兼容，行为同旧版）。

    Returns:
        导出的 ImageAsset 列表（已 MD5 去重）。
    """
    assets: List[ImageAsset] = []
    seen_hashes: dict = {}  # md5 -> 首次导出的 ImageAsset

    if document is None:
        try:
            import docx
        except ImportError:
            return []
        try:
            document = docx.Document(local_path)
        except Exception as e:
            print(f"      ⚠️ Failed to open DOCX for image extraction: {e}")
            return []

    doc_basename = os.path.splitext(os.path.basename(local_path))[0]
    img_index = 0

    for rel in document.part.rels.values():
        if "image" not in rel.reltype:
            continue

        try:
            blob = rel.target_part.blob
        except Exception:
            continue

        # MD5 去重：字节只导出一次，但不同 rId 指向同字节的"别名"关系
        # 必须保留各自的 target_ref，否则正文里第二处引用对齐不到资产、
        # 整个出现位置丢失（图② 同时出现在步骤1/步骤3 的场景）。
        md5 = hashlib.md5(blob).hexdigest()
        first = seen_hashes.get(md5)
        if first is not None:
            if rel.target_ref != first.original_name:
                assets.append(ImageAsset(
                    local_path=first.local_path,
                    page_num=None,
                    image_index=img_index,
                    original_name=rel.target_ref,
                ))
                img_index += 1
            continue

        # 确定文件扩展名
        content_type = getattr(rel.target_part, 'content_type', '')
        ext = _content_type_to_ext(content_type, rel.target_ref)

        # 跳过不支持的矢量格式（WMF/EMF/SVG）；G4：计数使损失可见
        if ext in (".wmf", ".emf", ".svg"):
            _note_skipped_vector(stats, ext)
            continue

        # 导出到 tmp_dir
        filename = f"{doc_basename}_img{img_index:04d}{ext}"
        out_path = os.path.join(output_dir, filename)

        try:
            with open(out_path, "wb") as f:
                f.write(blob)
        except Exception as e:
            print(f"      ⚠️ Failed to export DOCX image {rel.target_ref}: {e}")
            continue

        asset = ImageAsset(
            local_path=out_path,
            page_num=None,
            image_index=img_index,
            original_name=rel.target_ref,
        )
        assets.append(asset)
        seen_hashes[md5] = asset
        img_index += 1

    if assets:
        print(f"      [docx-img] Extracted {len(assets)} images "
              f"({len(seen_hashes)} unique blobs, "
              f"{len([r for r in document.part.rels.values() if 'image' in r.reltype])} refs)")

    return assets


# ═══════════════════════════════════════════════════════════════
# PDF — PyMuPDF (fitz) 逐页提取
# ═══════════════════════════════════════════════════════════════

# PIL transpose op → 日志名（lazy 引用 PIL，模块级只存 op 名）
_TRANSPOSE_LOG_NAMES = {
    "FLIP_TOP_BOTTOM": "flip_tb",
    "FLIP_LEFT_RIGHT": "flip_lr",
    "ROTATE_180": "rot180",
    "ROTATE_90": "rot90ccw",
    "ROTATE_270": "rot90cw",
    "TRANSPOSE": "transpose",
    "TRANSVERSE": "transverse",
}


def _pdf_page_display_ops(page) -> dict:
    """
    计算页面上每个图片 xref 的显示朝向校正操作。

    扫描件/转换器产出的 PDF 常把位图按非常规朝向存储（如逐行倒序＝垂直镜像、
    或旋转 180°），再用内容流 CTM 补偿，viewer 显示正常；但 extract_image
    导出的是存储字节，朝向就是错的。page /Rotate 同理只作用于显示。

    判定依据：get_image_info 的 transform（未旋转页坐标系）∘ page.rotation_matrix
    的符号模式 → 8 类朝向。轴主导时看 (a,d) 符号，对角主导时看 (b,c) 符号。
    （已用 9 种合成组合 + 真实坏档逐像素对照页面渲染验证。）

    Returns:
        {xref: PIL Transpose op 名称字符串}；缺失或 None = 无需校正。
    """
    ops = {}
    try:
        infos = page.get_image_info(xrefs=True)
    except Exception:
        return ops

    import fitz

    best_area = {}  # xref -> bbox 面积（同一 xref 多次出现时取最大实例）
    for info in infos:
        xref = info.get("xref", 0)
        transform = info.get("transform")
        if not xref or transform is None:
            continue
        bbox = info.get("bbox") or (0, 0, 0, 0)
        area = abs((bbox[2] - bbox[0]) * (bbox[3] - bbox[1]))
        if xref in best_area and area <= best_area[xref]:
            continue
        best_area[xref] = area

        try:
            mat = fitz.Matrix(transform) * page.rotation_matrix
        except Exception:
            continue
        a, b, c, d = mat.a, mat.b, mat.c, mat.d

        if abs(a) + abs(d) >= abs(b) + abs(c):
            # 轴对齐：a/d 符号决定上下/左右镜像
            key = (a >= 0, d >= 0)
            op = {
                (True, True): None,
                (True, False): "FLIP_TOP_BOTTOM",
                (False, True): "FLIP_LEFT_RIGHT",
                (False, False): "ROTATE_180",
            }[key]
        else:
            # 90° 族：b/c 符号决定旋转方向/转置
            key = (b >= 0, c >= 0)
            op = {
                (True, False): "ROTATE_270",   # 显示为顺时针 90°
                (False, True): "ROTATE_90",    # 显示为逆时针 90°
                (True, True): "TRANSPOSE",
                (False, False): "TRANSVERSE",
            }[key]

        ops[xref] = op
    return ops


def _pdf_page_image_bboxes(page) -> dict:
    """每个图片 xref 在本页的显示 bbox（同一 xref 多实例时取最大面积者）。

    与 _pdf_page_display_ops 同源（get_image_info），坐标为页空间、上原点 —
    与 pdfplumber 文本行的 top/bottom 同坐标系，供版面位置图片锚定。
    """
    bboxes = {}
    try:
        infos = page.get_image_info(xrefs=True)
    except Exception:
        return bboxes
    best_area = {}
    for info in infos:
        xref = info.get("xref", 0)
        bbox = info.get("bbox")
        if not xref or not bbox:
            continue
        area = abs((bbox[2] - bbox[0]) * (bbox[3] - bbox[1]))
        if xref not in best_area or area > best_area[xref]:
            best_area[xref] = area
            bboxes[xref] = tuple(float(v) for v in bbox)
    return bboxes


def _transpose_image_bytes(blob: bytes, op_name: str) -> Optional[bytes]:
    """按 op_name 转置图片字节，使其与页面显示朝向一致。

    失败返回 None（调用方回退写原始字节，保持优雅降级）。
    JPEG 需重编码（quality=95，保留 ICC）；PNG 等无损格式无损转置。
    """
    try:
        import io
        from PIL import Image
    except ImportError:
        return None

    try:
        transpose_enum = getattr(Image, "Transpose", Image)
        op = getattr(transpose_enum, op_name)
        src = Image.open(io.BytesIO(blob))
        fmt = src.format or "PNG"
        out_img = src.transpose(op)
        buf = io.BytesIO()
        save_kwargs = {}
        if fmt == "JPEG":
            save_kwargs["quality"] = 95
        icc = src.info.get("icc_profile")
        if icc:
            save_kwargs["icc_profile"] = icc
        out_img.save(buf, format=fmt, **save_kwargs)
        return buf.getvalue()
    except Exception:
        return None


def extract_images_from_pdf(
    local_path: str,
    output_dir: str,
    max_pages: int = 20,
    stats: Optional[dict] = None,
) -> List[ImageAsset]:
    """
    从 PDF 文件中提取嵌入图片，带精确 page_num。

    策略：使用 PyMuPDF (fitz) 逐页 get_images + extract_image。
    PyMuPDF 可提供每张图片所在的精确页码。

    朝向校正：extract_image 返回的是存储字节；若内容流 CTM / page /Rotate
    使显示朝向不同于存储朝向（典型：扫描件存储为垂直镜像或 180° 旋转），
    则按 _pdf_page_display_ops 的判定转置后再落盘，保证导出图与读者所见一致
    （下游 VLM 描述、OCR、钉钉卡片渲染都依赖这一点）。

    Args:
        local_path: PDF 文件的本地路径。
        output_dir: 图片导出目标目录。
        max_pages: 最大处理页数（与 pdf_extractor 保持一致）。

    Returns:
        导出的 ImageAsset 列表（已 MD5 去重）。
    """
    try:
        import fitz
    except ImportError:
        print("      ⚠️ PyMuPDF (fitz) not installed, skipping PDF image extraction")
        return []

    assets: List[ImageAsset] = []
    seen_hashes: set = set()

    try:
        pdf = fitz.open(local_path)
    except Exception as e:
        print(f"      ⚠️ Failed to open PDF for image extraction: {e}")
        return []

    doc_basename = os.path.splitext(os.path.basename(local_path))[0]
    img_index = 0
    corrected_counts: dict = {}  # op 名 -> 校正张数（用于日志）

    for page_idx in range(min(len(pdf), max_pages)):
        page = pdf[page_idx]
        page_num = page_idx + 1

        try:
            image_list = page.get_images(full=True)
        except Exception:
            continue

        display_ops = _pdf_page_display_ops(page) if image_list else {}
        page_bboxes = _pdf_page_image_bboxes(page) if image_list else {}

        for img_info in image_list:
            xref = img_info[0]

            try:
                base_image = pdf.extract_image(xref)
            except Exception:
                continue

            if not base_image or "image" not in base_image:
                continue

            blob = base_image["image"]
            op_name = display_ops.get(xref)

            # MD5 去重（PDF 中相同图片可能在多页出现）。
            # key 含校正 op：同一存储字节在不同页可能以不同朝向显示。
            md5 = hashlib.md5(blob).hexdigest()
            dedup_key = (md5, op_name)
            if dedup_key in seen_hashes:
                continue
            seen_hashes.add(dedup_key)

            ext = f".{base_image.get('ext', 'png')}"
            # 跳过不支持的格式；G4：计数使损失可见
            if ext in (".wmf", ".emf", ".svg"):
                _note_skipped_vector(stats, ext)
                continue

            # 朝向校正：转置失败回退原始字节（不阻断提取）
            if op_name is not None:
                corrected = _transpose_image_bytes(blob, op_name)
                if corrected is not None:
                    blob = corrected
                    log_name = _TRANSPOSE_LOG_NAMES.get(op_name, op_name)
                    corrected_counts[log_name] = corrected_counts.get(log_name, 0) + 1
                else:
                    print(f"      ⚠️ Orientation fix failed for xref={xref} "
                          f"(op={op_name}), keeping raw bytes")

            filename = f"{doc_basename}_p{page_num}_img{img_index:04d}{ext}"
            out_path = os.path.join(output_dir, filename)

            try:
                with open(out_path, "wb") as f:
                    f.write(blob)
            except Exception as e:
                print(f"      ⚠️ Failed to export PDF image xref={xref}: {e}")
                continue

            assets.append(ImageAsset(
                local_path=out_path,
                page_num=page_num,
                image_index=img_index,
                original_name=f"xref_{xref}",
                bbox=page_bboxes.get(xref),
            ))
            img_index += 1

    total_pages = min(len(pdf), max_pages)
    pdf.close()

    if assets:
        print(f"      [pdf-img] Extracted {len(assets)} unique images from {total_pages} pages")
        if corrected_counts:
            detail = ", ".join(f"{k}×{v}" for k, v in sorted(corrected_counts.items()))
            print(f"      [pdf-img] ⤿ orientation-corrected {sum(corrected_counts.values())} "
                  f"images ({detail})")

    return assets


# ═══════════════════════════════════════════════════════════════
# XLSX — openpyxl _images 遍历
# ═══════════════════════════════════════════════════════════════

def hidden_sheet_should_ingest(ws) -> bool:
    """B2-1 条件化（release-gate 2026-07-06 复盘）：隐藏 sheet 是否入库。

    B2-1 无条件跳过隐藏 sheet 的初衷（草稿/中间计算/下拉源=噪声+泄露面）成立，但
    xlsx_clean 设备清扫基准书证明**作者藏 tab ≠ 内容作废**：2/3 的 sheet（30 行×15 列
    清扫基准表、13/24 张锚定图片）是隐藏的正文主体，无条件跳过使 binding Jaccard
    0.8917→0.6。本函数改为证据制：

      · visible          → 恒入库（与历史一致）
      · veryHidden 等    → 恒跳过（程序级深藏 = 明确的"不给人看"意图）
      · hidden           → 有**正文证据**才入库：带锚定图片（草稿/下拉源几乎不含图），
                           或 ≥10 非空行 × ≥3 列的二维内容体量（下拉源多为 1-2 列窄表）

    RAG_XLSX_HIDDEN_SHEET_MODE：auto（默认，证据制）| skip（B2-1 原行为，全跳）|
    ingest（B2-1 前行为，全收）。文本路径与图片路径共用本函数保持决策一致；
    read_only workbook 无 _images（大表 >50k 行/100MB 才触发）→ 图片证据不可见，
    但该体量下内容体量 prong 必然命中，两路径决策仍一致。任何读取异常 → 保守跳过。
    """
    state = getattr(ws, "sheet_state", "visible")
    if state == "visible":
        return True
    mode = os.environ.get("RAG_XLSX_HIDDEN_SHEET_MODE", "auto").strip().lower()
    if mode == "ingest":
        return True
    if mode == "skip":
        return False
    if state != "hidden":          # veryHidden / 未知态：恒跳过
        return False
    try:
        if getattr(ws, "_images", None):
            return True
        nonempty = 0
        for row in ws.iter_rows(max_row=min(ws.max_row or 0, 200)):
            if any(c.value not in (None, "") for c in row):
                nonempty += 1
                if nonempty >= 10:
                    break
        return nonempty >= 10 and (ws.max_column or 0) >= 3
    except Exception:
        return False               # 读取失败 → 维持 B2-1 保守跳过


def extract_images_from_xlsx(
    local_path: str,
    output_dir: str,
    workbook=None,
) -> List[ImageAsset]:
    """
    从 XLSX 文件中提取嵌入图片。

    策略：遍历 openpyxl 每个 worksheet 的 _images 列表。

    Args:
        local_path: XLSX 文件的本地路径。
        output_dir: 图片导出目标目录。
        workbook: 可选，调用方已加载的**普通模式** openpyxl workbook（perf#62 复用，
            避免同一 xlsx 再次 zip 解包）。read_only workbook 的 worksheet 没有
            _images（openpyxl 只在普通模式加载 drawing），直接用会静默丢图 →
            检测到 read_only 一律忽略、回退自行加载。传入的 workbook 不在此关闭
            （生命周期归调用方）；普通模式 _images 数据在 load 时已读入内存，
            不依赖 archive 开闭。

    Returns:
        导出的 ImageAsset 列表（已 MD5 去重）。
    """
    try:
        from openpyxl import load_workbook
        from openpyxl.drawing.image import Image as XlImage  # noqa: F401 — 可用性探测（缺则优雅降级返回 []）
    except ImportError:
        return []

    assets: List[ImageAsset] = []
    seen_hashes: dict = {}  # md5 -> {"asset": 首次导出资产, "anchors": set((sheet_idx, anchor_row))}

    wb = None
    owns_wb = True
    if workbook is not None and not getattr(workbook, "read_only", False):
        wb = workbook
        owns_wb = False
    if wb is None:
        try:
            wb = load_workbook(local_path, data_only=True)
        except Exception as e:
            print(f"      ⚠️ Failed to open XLSX for image extraction: {e}")
            return []

    doc_basename = os.path.splitext(os.path.basename(local_path))[0]
    img_index = 0

    for sheet_idx, ws in enumerate(wb.worksheets):
        # 隐藏 sheet 图片是否提取与文本路径共用 hidden_sheet_should_ingest 证据制决策
        # （B2-1 条件化）。enumerate 不跳号 → sheet_idx 与文本路径对齐
        # （drawing→sheet 的 MD5 字节匹配仍兜底）。
        if not hidden_sheet_should_ingest(ws):
            continue
        if not hasattr(ws, '_images'):
            continue

        for xl_img in ws._images:
            try:
                # openpyxl Image 对象的 _data 方法或 ref 属性
                if hasattr(xl_img, '_data') and callable(xl_img._data):
                    blob = xl_img._data()
                elif hasattr(xl_img, 'ref') and hasattr(xl_img.ref, 'read'):
                    xl_img.ref.seek(0)
                    blob = xl_img.ref.read()
                else:
                    continue
            except Exception:
                continue

            if not blob:
                continue

            md5 = hashlib.md5(blob).hexdigest()

            # 先取 anchor 行号（行级图片绑定 + 别名判定都要用）
            anchor_row = None
            anchor = getattr(xl_img, 'anchor', None)
            if anchor and hasattr(anchor, '_from'):
                anchor_row = getattr(anchor._from, 'row', None)

            # MD5 去重：字节只落盘一次；但同一张图复用在【不同 anchor 行】（步骤1/步骤3 共用
            # 一张截图）必须各自保留 anchor_row，否则第二处 step_card 丢图（与 DOCX 别名逻辑对齐）。
            dup = seen_hashes.get(md5)
            if dup is not None:
                key = (sheet_idx, anchor_row)
                if anchor_row is not None and key not in dup["anchors"]:
                    alias = ImageAsset(
                        local_path=dup["asset"].local_path,  # 复用已落盘文件，不重复写
                        page_num=sheet_idx + 1,
                        image_index=img_index,
                        original_name=f"sheet{sheet_idx}_image",
                        md5=md5,  # perf#88：别名与原资产同字节，直接携带同一 MD5
                    )
                    alias.anchor_row = anchor_row
                    assets.append(alias)
                    dup["anchors"].add(key)
                    img_index += 1
                continue

            # 推测扩展名
            ext = ".png"  # openpyxl 默认
            if blob[:3] == b'\xff\xd8\xff':
                ext = ".jpeg"
            elif blob[:4] == b'\x89PNG':
                ext = ".png"

            filename = f"{doc_basename}_sheet{sheet_idx}_img{img_index:04d}{ext}"
            out_path = os.path.join(output_dir, filename)

            try:
                with open(out_path, "wb") as f:
                    f.write(blob)
            except Exception as e:
                print(f"      ⚠️ Failed to export XLSX image: {e}")
                continue

            asset = ImageAsset(
                local_path=out_path,
                page_num=sheet_idx + 1,
                image_index=img_index,
                original_name=f"sheet{sheet_idx}_image",
                md5=md5,  # perf#88：导出时已算过的 MD5 直接携带，配对阶段免重读重哈希
            )
            if anchor_row is not None:
                asset.anchor_row = anchor_row
            assets.append(asset)
            seen_hashes[md5] = {"asset": asset, "anchors": {(sheet_idx, anchor_row)}}
            img_index += 1

    if owns_wb:
        wb.close()

    # ── 后处理：从 Drawing XML 提取 group 编号标注 ──
    # openpyxl 的 _images 无法读取 grpSp 里的文字标注（如①②③序号圆圈）
    # 需要直接解析 xlsx zip 内的 drawing XML
    if assets:
        try:
            _enrich_xlsx_annotations(local_path, assets, doc_basename)
        except Exception as e:
            print(f"      ⚠️ Drawing XML annotation extraction failed: {e}")

        print(f"      [xlsx-img] Extracted {len(assets)} unique images from {len(wb.worksheets)} sheets")

    return assets


def _enrich_xlsx_annotations(xlsx_path: str, assets: List[ImageAsset], doc_basename: str):
    """从 XLSX Drawing XML 解析标注编号，绑定到对应图片。

    支持两种 Drawing 结构：
      A) grpSp 分组：图片和标注在同一个 <xdr:grpSp> 内（直接配对）
      B) standalone：图片和标注是独立的 <xdr:twoCellAnchor>（按坐标邻近配对）
    """
    import zipfile
    import xml.etree.ElementTree as ET

    ns = {
        'xdr': 'http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }

    def _get_anchor_pos(anchor_el):
        """从 <xdr:from> 提取 (row, col, colOff) 绝对位置。"""
        from_el = anchor_el.find('xdr:from', ns)
        if from_el is None:
            return None, None, 0
        row = int(from_el.find('xdr:row', ns).text)
        col = int(from_el.find('xdr:col', ns).text)
        col_off = int(from_el.find('xdr:colOff', ns).text)
        return row, col, col_off

    # perf#88：导出循环已算过每个 blob 的 MD5（ImageAsset.md5），此处一次性预建
    # {md5: [assets 原序]} 索引——方式A/方式B 每次配对从「全量重读文件+重哈希」
    # （O(pics×assets) 磁盘 IO）降为零额外 IO 的字典查找。asset.md5 为 None
    # （外部构造的旧路径，如 docx_extractor）时回退现状逐文件读哈希，每文件至多
    # 一次；读失败按现状跳过该资产。
    md5_to_assets: dict = {}
    for asset in assets:
        asset_md5 = asset.md5
        if asset_md5 is None:
            try:
                with open(asset.local_path, 'rb') as f:
                    asset_md5 = hashlib.md5(f.read()).hexdigest()
            except Exception:
                continue
        md5_to_assets.setdefault(asset_md5, []).append(asset)

    def _find_asset_by_md5(media_blob, sheet_idx):
        """通过 MD5 找到对应的 ImageAsset（预建索引；assets 原序 first-match 语义不变）。"""
        if not media_blob:
            return None
        media_md5 = hashlib.md5(media_blob).hexdigest()
        for asset in md5_to_assets.get(media_md5, ()):
            if asset.page_num == sheet_idx + 1:
                return asset
        return None

    def _drawing_to_sheet_order(z):
        """批次6（ultra image_extraction_utils:672）：解析 drawing 文件 → workbook 顺序
        sheet_idx（0-based）的**真实**映射。OOXML 的 drawingN.xml 按「有图 sheet 的创建序」
        编号，不是按 sheet 序——首 sheet 无图（封面/目录页）时 sheet2 的 drawing 是
        drawing1.xml，旧的 drawing_num-1 假定令 sheet_idx=0 而资产 page_num=2，
        _find_asset_by_md5 的 page_num 闸全数落空 → 标注绑定静默全丢（①②③ 显式绑定
        降级为弱位置启发）。真实链条：workbook.xml 的 <sheet> 顺序（page_num 与
        openpyxl wb.worksheets 同源）→ r:id → workbook.xml.rels → worksheets/sheetN.xml
        → 其 _rels → ../drawings/drawingM.xml。任何一环缺失/解析失败 → 返回 {}，
        调用方回退旧假定（fail-open，常规文件两者一致时零行为变化）。"""
        _R_NS = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        _MAIN_NS = {'x': 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'}
        mapping = {}
        names = set(z.namelist())
        wb_root = ET.fromstring(z.read('xl/workbook.xml'))
        rels_root = ET.fromstring(z.read('xl/_rels/workbook.xml.rels'))
        rid_to_target = {rel.get('Id'): rel.get('Target', '') or '' for rel in rels_root}
        sheets = wb_root.findall('.//x:sheets/x:sheet', _MAIN_NS)
        for order_idx, sheet_el in enumerate(sheets):
            target = rid_to_target.get(sheet_el.get(_R_NS), '')
            if not target:
                continue
            sheet_file = target.lstrip('/').split('/')[-1]        # sheet3.xml
            ws_rels = f'xl/worksheets/_rels/{sheet_file}.rels'
            if ws_rels not in names:
                continue
            try:
                ws_rels_root = ET.fromstring(z.read(ws_rels))
            except ET.ParseError:
                continue
            for rel in ws_rels_root:
                m2 = re.search(r'drawing(\d+)\.xml$', rel.get('Target', '') or '')
                if m2:
                    mapping[f'xl/drawings/drawing{m2.group(1)}.xml'] = order_idx
        return mapping

    import re
    with zipfile.ZipFile(xlsx_path) as z:
        drawing_files = sorted([n for n in z.namelist()
                                if n.startswith('xl/drawings/drawing') and n.endswith('.xml')])
        rels_files = {n for n in z.namelist()
                      if n.startswith('xl/drawings/_rels/') and n.endswith('.rels')}
        try:
            _sheet_order_map = _drawing_to_sheet_order(z)
        except Exception as _me:
            print(f"      ⚠️ [xlsx-img] drawing↔sheet rels 解析失败，回退 drawingN→sheetN 假定: {_me}")
            _sheet_order_map = {}

        for drawing_path in drawing_files:
            m = re.search(r'drawing(\d+)\.xml$', drawing_path)
            if not m:
                continue
            drawing_num = int(m.group(1))
            if drawing_path in _sheet_order_map:
                sheet_idx = _sheet_order_map[drawing_path]
            else:
                sheet_idx = drawing_num - 1   # 回退旧假定（rels 链解析不出该 drawing）

            # 解析 rels → rId → media filename
            rels_path = f'xl/drawings/_rels/drawing{drawing_num}.xml.rels'
            rid_to_media = {}
            if rels_path in rels_files:
                rels_root = ET.fromstring(z.read(rels_path))
                for rel in rels_root:
                    rid = rel.get('Id', '')
                    target = rel.get('Target', '')
                    if 'image' in target.lower() or 'media' in target.lower():
                        rid_to_media[rid] = os.path.basename(target)

            if not rid_to_media:
                continue

            root = ET.fromstring(z.read(drawing_path))

            # ── 方式 A: grpSp 分组配对 ──
            grp_matched = set()  # 已通过 grpSp 配对的 asset id
            for anchor in root.findall('xdr:twoCellAnchor', ns):
                grp = anchor.find('xdr:grpSp', ns)
                if grp is None:
                    continue

                pics = grp.findall('.//xdr:pic', ns)
                sps = grp.findall('.//xdr:sp', ns)
                if not pics:
                    continue

                group_texts = []
                for sp in sps:
                    for t_el in sp.findall('.//a:t', ns):
                        if t_el.text and t_el.text.strip():
                            group_texts.append(t_el.text.strip())

                annotation_nums = [t for t in group_texts if t.isdigit()]
                annotation_labels = [t for t in group_texts if not t.isdigit()]

                for pic in pics:
                    blip = pic.find(
                        './/{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                    if blip is None:
                        continue
                    r_id = blip.get(
                        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', '')
                    if not r_id or r_id not in rid_to_media:
                        continue

                    media_name = rid_to_media[r_id]
                    media_path = f'xl/media/{media_name}'
                    media_blob = z.read(media_path) if media_path in z.namelist() else None
                    matched_asset = _find_asset_by_md5(media_blob, sheet_idx)
                    if matched_asset:
                        if annotation_nums:
                            matched_asset.annotation_num = int(annotation_nums[0])
                        if annotation_labels:
                            matched_asset.annotation_label = ', '.join(annotation_labels)
                        grp_matched.add(id(matched_asset))

            # ── 方式 B: standalone 邻近配对 ──
            # 收集所有独立的 PIC 和 TEXT anchor
            standalone_pics = []  # [(row, col, col_off, rId)]
            standalone_texts = []  # [(row, col, col_off, text)]

            for anchor in root.findall('xdr:twoCellAnchor', ns):
                if anchor.find('xdr:grpSp', ns) is not None:
                    continue  # 已在方式 A 处理

                row, col, col_off = _get_anchor_pos(anchor)
                if row is None:
                    continue

                pic = anchor.find('xdr:pic', ns)
                sp = anchor.find('xdr:sp', ns)

                if pic is not None:
                    blip = pic.find(
                        './/{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                    if blip is not None:
                        r_id = blip.get(
                            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', '')
                        if r_id and r_id in rid_to_media:
                            standalone_pics.append((row, col, col_off, r_id))
                elif sp is not None:
                    texts = [t.text for t in sp.findall('.//a:t', ns) if t.text]
                    text = ''.join(texts).strip()
                    if text and text.isdigit():
                        standalone_texts.append((row, col, col_off, text))

            if standalone_pics and standalone_texts:
                # 全局贪心匹配：计算所有 (TEXT, PIC) 对的距离，按距离排序，
                # 每个 PIC 和 TEXT 只匹配一次
                EMU_PER_COL = 9525 * 100

                all_pairs = []  # [(distance, text_idx, pic_idx)]
                for ti, (t_row, t_col, t_col_off, t_text) in enumerate(standalone_texts):
                    t_abs = t_col * EMU_PER_COL + t_col_off
                    for pi, (p_row, p_col, p_col_off, p_rid) in enumerate(standalone_pics):
                        row_dist = abs(p_row - t_row)
                        if row_dist > 5:
                            continue
                        p_abs = p_col * EMU_PER_COL + p_col_off
                        col_dist = abs(p_abs - t_abs)
                        total_dist = row_dist * 100000000 + col_dist
                        all_pairs.append((total_dist, ti, pi))

                all_pairs.sort()  # 按距离从近到远

                matched_texts = set()
                matched_pics = set()

                for _, ti, pi in all_pairs:
                    if ti in matched_texts or pi in matched_pics:
                        continue  # 已配对，跳过

                    t_text = standalone_texts[ti][3]
                    p_rid = standalone_pics[pi][3]

                    media_name = rid_to_media[p_rid]
                    media_path = f'xl/media/{media_name}'
                    media_blob = z.read(media_path) if media_path in z.namelist() else None
                    matched_asset = _find_asset_by_md5(media_blob, sheet_idx)
                    if matched_asset and id(matched_asset) not in grp_matched:
                        matched_asset.annotation_num = int(t_text)
                        matched_texts.add(ti)
                        matched_pics.add(pi)


# ═══════════════════════════════════════════════════════════════
# PPTX — python-pptx relationship 遍历
# ═══════════════════════════════════════════════════════════════

def extract_images_from_pptx(
    local_path: str,
    output_dir: str,
    stats: Optional[dict] = None,
) -> List[ImageAsset]:
    """
    从 PPTX 文件中提取所有嵌入图片。

    策略：PPTX 是 OOXML 格式，和 DOCX 类似。
    遍历每个 slide 的 part.rels 中 reltype 包含 'image' 的关系，
    读取 target_part.blob 获取图片二进制数据。

    每张图片会记录所在的 slide 编号作为 page_num。

    Args:
        local_path: PPTX 文件的本地路径。
        output_dir: 图片导出目标目录。

    Returns:
        导出的 ImageAsset 列表（已 MD5 去重）。
    """
    try:
        from pptx import Presentation
    except ImportError:
        return []

    assets: List[ImageAsset] = []
    seen_hashes: set = set()

    try:
        prs = Presentation(local_path)
    except Exception as e:
        print(f"      ⚠️ Failed to open PPTX for image extraction: {e}")
        return []

    doc_basename = os.path.splitext(os.path.basename(local_path))[0]
    img_index = 0

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1

        for rel in slide.part.rels.values():
            if "image" not in rel.reltype:
                continue

            try:
                blob = rel.target_part.blob
            except Exception:
                continue

            # MD5 去重
            md5 = hashlib.md5(blob).hexdigest()
            if md5 in seen_hashes:
                continue
            seen_hashes.add(md5)

            # 确定文件扩展名
            content_type = getattr(rel.target_part, 'content_type', '')
            ext = _content_type_to_ext(content_type, rel.target_ref)

            # 跳过不支持的矢量格式（WMF/EMF/SVG）；G4：计数使损失可见
            if ext in (".wmf", ".emf", ".svg"):
                _note_skipped_vector(stats, ext)
                continue

            # 导出到 tmp_dir
            filename = f"{doc_basename}_slide{slide_num}_img{img_index:04d}{ext}"
            out_path = os.path.join(output_dir, filename)

            try:
                with open(out_path, "wb") as f:
                    f.write(blob)
            except Exception as e:
                print(f"      ⚠️ Failed to export PPTX image {rel.target_ref}: {e}")
                continue

            assets.append(ImageAsset(
                local_path=out_path,
                page_num=slide_num,
                image_index=img_index,
                original_name=rel.target_ref,
            ))
            img_index += 1

    total_rels = sum(
        1 for slide in prs.slides
        for rel in slide.part.rels.values()
        if "image" in rel.reltype
    )
    if assets:
        print(f"      [pptx-img] Extracted {len(assets)} unique images "
              f"(deduped from {total_rels} refs across {len(prs.slides)} slides)")

    return assets


# ═══════════════════════════════════════════════════════════════
# 辅助函数
# ═══════════════════════════════════════════════════════════════

def _content_type_to_ext(content_type: str, target_ref: str = "") -> str:
    """将 MIME content_type 或文件引用名转换为文件扩展名。"""
    ct_map = {
        "image/jpeg": ".jpeg",
        "image/jpg": ".jpeg",
        "image/png": ".png",
        "image/gif": ".gif",
        "image/bmp": ".bmp",
        "image/tiff": ".tiff",
        "image/webp": ".webp",
        "image/x-wmf": ".wmf",
        "image/x-emf": ".emf",
        "image/svg+xml": ".svg",
    }
    if content_type in ct_map:
        return ct_map[content_type]

    # fallback: 从 target_ref 推断（如 media/image3.jpeg）
    if target_ref:
        _, ext = os.path.splitext(target_ref)
        if ext:
            return ext.lower()

    return ".png"
