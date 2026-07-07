# -*- coding: utf-8 -*-
"""tests/test_tier1_extraction_fixes.py — 工业级审计第一梯队修复的回归测试。

G4:  WMF/EMF/SVG 矢量图跳过必须计数并落 [VECTOR_IMAGE_SKIPPED] warning（此前静默丢弃）。
G13: PPTX GroupShape 递归——组合形状内文本框此前静默消失。
G14: 加密/带密码文件前置探测——此前只表现为通用打开失败并空烧 OCR 与重试。
G15: 乱码质量门——坏字体页 ≥50 字符 mojibake 此前永不重 OCR。
"""
import os
import sys
import types


sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
os.environ.setdefault("RAG_ENV", "test")


# ═══════════════════ G4: 矢量图跳过计数 ═══════════════════

class _FakePart:
    def __init__(self, blob, content_type):
        self.blob = blob
        self.content_type = content_type


class _FakeRel:
    def __init__(self, target_ref, blob, content_type):
        self.reltype = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
        self.target_ref = target_ref
        self.target_part = _FakePart(blob, content_type)


def _fake_document(rels):
    part = types.SimpleNamespace(rels={f"rId{i}": r for i, r in enumerate(rels)})
    return types.SimpleNamespace(part=part)


def test_g4_docx_emf_skip_is_counted(tmp_path):
    from opensearch_pipeline.extraction.image_extraction_utils import extract_images_from_docx
    doc = _fake_document([
        _FakeRel("media/image1.emf", b"\x01\x00\x00\x00emf-bytes", "image/x-emf"),
        _FakeRel("media/image2.wmf", b"\xd7\xcd\xc6\x9awmf-bytes", "image/x-wmf"),
        _FakeRel("media/image3.png", b"\x89PNG\r\n\x1a\npng-bytes", "image/png"),
    ])
    stats = {}
    assets = extract_images_from_docx("/tmp/fake.docx", str(tmp_path), document=doc, stats=stats)
    assert stats["skipped_vector_images"] == 2
    assert stats["skipped_vector_formats"] == {"emf", "wmf"}
    assert len(assets) == 1 and assets[0].local_path.endswith(".png")


def test_g4_stats_none_keeps_old_behavior(tmp_path):
    """独立调用（不传 stats）不报错、跳过行为不变。"""
    from opensearch_pipeline.extraction.image_extraction_utils import extract_images_from_docx
    doc = _fake_document([_FakeRel("media/a.svg", b"<svg/>", "image/svg+xml")])
    assets = extract_images_from_docx("/tmp/fake.docx", str(tmp_path), document=doc)
    assert assets == []


# ═══════════════════ G13: PPTX GroupShape 递归 ═══════════════════

def test_g13_group_shape_text_extracted(tmp_path):
    """组合形状（含嵌套组）内的文本框此前被 hasattr(shape,'text') 静默跳过。"""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(914400), Emu(457200))
    tb.text_frame.text = "顶层文本框"
    group = slide.shapes.add_group_shape()
    g_tb = group.shapes.add_textbox(Emu(0), Emu(500000), Emu(914400), Emu(457200))
    g_tb.text_frame.text = "组内说明文字"
    nested = group.shapes.add_group_shape()
    n_tb = nested.shapes.add_textbox(Emu(0), Emu(1000000), Emu(914400), Emu(457200))
    n_tb.text_frame.text = "嵌套组内文字"
    path = str(tmp_path / "g13.pptx")
    prs.save(path)

    from opensearch_pipeline.extraction.unified_extractor import UnifiedExtractor
    res = UnifiedExtractor(simulate=True).extract({
        "doc_id": "g13", "version_no": 1, "file_ext": "pptx", "raw_key": "raw/g13.pptx",
        "filename": "g13.pptx", "local_path": path, "_tmp_dir": str(tmp_path),
    })
    assert "顶层文本框" in res.text
    assert "组内说明文字" in res.text      # the bug: previously silently dropped
    assert "嵌套组内文字" in res.text      # recursion depth > 1


# ═══════════════════ G15: 乱码质量门 ═══════════════════

def test_g15_quality_predicate():
    from opensearch_pipeline.extraction.unified_extractor import _native_text_quality_ok
    # 正常中文/中英混排/表格数字 → OK
    assert _native_text_quality_ok("一、审核流程：核对发票信息，金额 ¥1,234.56（含税）。" * 3)
    assert _native_text_quality_ok("Step 1: open U8+ ERP system, click 采购管理 -> 订单。" * 3)
    # (cid:N) 串 → 坏字体
    assert not _native_text_quality_ok("(cid:123)(cid:456)(cid:789)" * 10)
    # PUA 乱码 → 坏字体
    assert not _native_text_quality_ok("".join(chr(0xE000 + i % 100) for i in range(80)))
    # 替换符密度超阈
    assert not _native_text_quality_ok(("正常文字" + "�" * 3) * 20)
    # 空文本不触发（交给字符数门槛）
    assert _native_text_quality_ok("   ")


def _pdf_result(per_page_texts, page_count):
    from opensearch_pipeline.extraction.schema import ExtractedBlock, ExtractionResult
    blocks = [ExtractedBlock(block_type="paragraph", text=t, page_num=pg)
              for pg, t in per_page_texts.items()]
    text = "\n".join(per_page_texts.values())
    return ExtractionResult(doc_id="d", version_no=1, source_key="", file_ext="pdf",
                            extract_method="m", title="t", text=text,
                            text_length=len(text), blocks=blocks, page_count=page_count)


def test_g15_garbled_page_flagged_for_ocr():
    from opensearch_pipeline.extraction.unified_extractor import UnifiedExtractor
    ux = UnifiedExtractor.__new__(UnifiedExtractor)
    res = _pdf_result({1: "正常的中文正文内容，" * 10,
                       2: "(cid:12)(cid:34)(cid:56)" * 10}, page_count=2)  # ≥50 字符的乱码页
    assert ux._pages_needing_ocr(res) == [2]
    assert res._garbled_pages == {2}
    assert any("[GARBLED_TEXT]" in w for w in res.warnings)


def test_g15_clean_pages_unaffected():
    from opensearch_pipeline.extraction.unified_extractor import UnifiedExtractor
    ux = UnifiedExtractor.__new__(UnifiedExtractor)
    res = _pdf_result({1: "正常的中文正文内容，" * 10, 2: "second page body text " * 5},
                      page_count=2)
    assert ux._pages_needing_ocr(res) == []
    assert not hasattr(res, "_garbled_pages")


def test_g15_ocr_replaces_garbled_native_blocks(monkeypatch):
    """OCR 成功后乱码页的原生块被剔除（不与 OCR 文本并存）；OCR 未覆盖的页保留。"""
    from opensearch_pipeline.extraction.schema import ExtractedBlock
    from opensearch_pipeline.extraction.unified_extractor import UnifiedExtractor
    ux = UnifiedExtractor.__new__(UnifiedExtractor)
    res = _pdf_result({1: "正常的中文正文内容，" * 10,
                       2: "(cid:12)(cid:34)(cid:56)" * 10}, page_count=2)
    needy = ux._pages_needing_ocr(res)
    assert needy == [2]

    ocr_block = ExtractedBlock(block_type="ocr_text", text="OCR 重取的第二页正文",
                               page_num=2, source="ocr")
    fake_ocr = types.SimpleNamespace(
        to_blocks=lambda: [ocr_block], status="SUCCESS", combined_text="OCR 重取的第二页正文")
    ux.ocr_client = types.SimpleNamespace(ocr_pdf=lambda *a, **k: fake_ocr)
    ux.oss_client = None

    out = ux._apply_ocr_fallback({"doc_id": "d", "local_path": "/tmp/x.pdf"}, res, needy=needy)
    assert "OCR 重取的第二页正文" in out.text
    assert "(cid:12)" not in out.text                    # 乱码原生块已剔除
    assert "正常的中文正文内容" in out.text               # 干净页不受影响


# ═══════════════════ G14: 加密文件前置探测 ═══════════════════

def _extract_task(path, ext, tmp_path):
    return {"doc_id": "g14", "version_no": 1, "file_ext": ext, "raw_key": f"raw/g14.{ext}",
            "filename": f"g14.{ext}", "local_path": path, "_tmp_dir": str(tmp_path)}


def test_g14_encrypted_pdf_short_circuits(tmp_path):
    from pypdf import PdfReader, PdfWriter
    import fitz
    src = fitz.open()
    src.new_page().insert_text((72, 72), "secret content")
    plain = str(tmp_path / "plain.pdf")
    src.save(plain); src.close()
    writer = PdfWriter(clone_from=plain)
    writer.encrypt(user_password="s3cret")
    enc = str(tmp_path / "enc.pdf")
    with open(enc, "wb") as f:
        writer.write(f)
    assert PdfReader(enc).is_encrypted

    from opensearch_pipeline.extraction.unified_extractor import UnifiedExtractor
    ux = UnifiedExtractor(simulate=True)
    res = ux.extract(_extract_task(enc, "pdf", tmp_path))
    assert res.extract_method == "unsupported:password_protected_pdf"
    assert res.ocr_status == "NOT_REQUIRED"          # 不空烧 OCR
    assert any("[PASSWORD_PROTECTED]" in w for w in res.warnings)


def test_g14_owner_password_only_pdf_passes_through(tmp_path):
    """仅 owner-password 限权（空用户密码、内容可读）的 PDF 应放行走正常抽取。"""
    from pypdf import PdfWriter
    import fitz
    src = fitz.open()
    src.new_page().insert_text((72, 72), "readable body text " * 10)
    plain = str(tmp_path / "p2.pdf")
    src.save(plain); src.close()
    writer = PdfWriter(clone_from=plain)
    writer.encrypt(user_password="", owner_password="admin")
    enc = str(tmp_path / "owner_only.pdf")
    with open(enc, "wb") as f:
        writer.write(f)

    from opensearch_pipeline.extraction.unified_extractor import _detect_password_protected
    assert _detect_password_protected(enc, "pdf") is None


def test_g14_ole_wrapped_docx_detected(tmp_path):
    """Office 加密 OOXML（或误改后缀的旧版二进制）文件头是 OLE 魔数而非 PK。"""
    fake = str(tmp_path / "enc.docx")
    with open(fake, "wb") as f:
        f.write(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 64)

    from opensearch_pipeline.extraction.unified_extractor import UnifiedExtractor
    ux = UnifiedExtractor(simulate=True)
    res = ux.extract(_extract_task(fake, "docx", tmp_path))
    assert res.extract_method == "unsupported:password_protected_docx"
    assert any("[PASSWORD_PROTECTED]" in w for w in res.warnings)


def test_g14_normal_docx_not_flagged(tmp_path):
    import docx
    d = docx.Document()
    d.add_paragraph("正常文档")
    path = str(tmp_path / "ok.docx")
    d.save(path)
    from opensearch_pipeline.extraction.unified_extractor import _detect_password_protected
    assert _detect_password_protected(path, "docx") is None


# ═══════════════════ G21: 图片 OCR 脱敏消费点兜底 ═══════════════════

def _canon_with_asset(ocr_text):
    body = "这是用于切块的正文内容，描述了操作步骤与注意事项。" * 30
    return {"doc_id": "D21", "version_no": 1, "file_ext": "pdf",
            "text": body, "title": "x", "category_l1": "", "category_l2": "",
            "blocks": [], "assets": [{"ocr_text": ocr_text, "oss_key": "k"}]}


def test_g21_failclosed_masks_pii_in_asset_ocr(monkeypatch):
    """RAG_IMAGE_OCR_PII（node 03）未开时，chunk 消费点兜底掩码 asset ocr_text。"""
    import opensearch_pipeline.pipeline_nodes as pn
    monkeypatch.delenv("RAG_IMAGE_OCR_PII", raising=False)
    monkeypatch.delenv("RAG_IMAGE_OCR_PII_FAILCLOSED", raising=False)  # 默认 ON
    doc = _canon_with_asset("联系人电话 13812345678，工号 A123")
    pn.node_chunk_documents({"canonicals": [doc]})
    assert "13812345678" not in doc["assets"][0]["ocr_text"]
    assert "138****5678" in doc["assets"][0]["ocr_text"]


def test_g21_material_code_fp_not_masked(monkeypatch):
    """物料编码锚点共现时，18 位编码命中 cn_id_card 属已知 FP，不掩码。"""
    import opensearch_pipeline.pipeline_nodes as pn
    monkeypatch.delenv("RAG_IMAGE_OCR_PII_FAILCLOSED", raising=False)
    fp_text = "物料编码 110101199003071234 规格说明"
    doc = _canon_with_asset(fp_text)
    pn.node_chunk_documents({"canonicals": [doc]})
    assert doc["assets"][0]["ocr_text"] == fp_text


def test_g21_failclosed_can_be_disabled(monkeypatch):
    import opensearch_pipeline.pipeline_nodes as pn
    monkeypatch.setenv("RAG_IMAGE_OCR_PII_FAILCLOSED", "false")
    raw = "电话 13812345678"
    doc = _canon_with_asset(raw)
    pn.node_chunk_documents({"canonicals": [doc]})
    assert doc["assets"][0]["ocr_text"] == raw


def test_g21_failclosed_masks_pii_in_visual_summary(monkeypatch):
    """106E77 事故回归：VLM caption（visual_summary）里的真手机号也要兜底掩码。"""
    import opensearch_pipeline.pipeline_nodes as pn
    monkeypatch.delenv("RAG_IMAGE_OCR_PII", raising=False)
    monkeypatch.delenv("RAG_IMAGE_OCR_PII_FAILCLOSED", raising=False)
    doc = _canon_with_asset("")
    doc["assets"][0]["visual_summary"] = "名片图片，联系电话 13812345678"
    pn.node_chunk_documents({"canonicals": [doc]})
    vs = doc["assets"][0]["visual_summary"]
    assert "13812345678" not in vs and "138****5678" in vs


def test_g21_credential_number_not_masked_in_caption(monkeypatch):
    """凭证锚点（注册号/证书）上下文的号码是业务事实，不得掩码——106E77 FDA 教训。"""
    import opensearch_pipeline.pipeline_nodes as pn
    monkeypatch.delenv("RAG_IMAGE_OCR_PII_FAILCLOSED", raising=False)
    cap = "FDA注册确认函，显示注册号13889211462、企业名称Fuling"
    doc = _canon_with_asset("")
    doc["assets"][0]["visual_summary"] = cap
    pn.node_chunk_documents({"canonicals": [doc]})
    assert doc["assets"][0]["visual_summary"] == cap


def test_scrub_image_text_unit():
    """scrub_image_text 直接单测：真手机掩码、凭证号保留、重复处理不泄露。"""
    from opensearch_pipeline.pii_patterns import scrub_image_text
    assert scrub_image_text("手机 13812345678") == "手机 138****5678"
    assert scrub_image_text("注册号 13889211462") == "注册号 13889211462"
    # 证书编号（17 位数字命中 bank_card）在凭证锚点上下文保留——106E77 教训
    assert scrub_image_text("certificate number is IDA00220120012001201") \
        == "certificate number is IDA00220120012001201"
    # 重复处理的安全契约：绝不重新暴露原始号码，且收敛（masked_id 会把 138****5678
    # 进一步吞成 [标识已脱敏]——牺牲可读性但不泄露，是既有 masked_id 语义，非 bug）。
    once = scrub_image_text("手机 13812345678")
    twice = scrub_image_text(once)
    assert "13812345678" not in twice
    assert scrub_image_text(twice) == twice  # 三次收敛


# ═══════════════════ G25: chunk 级乱码过滤 ═══════════════════

def test_g25_gibberish_predicate():
    import opensearch_pipeline.pipeline_nodes as pn
    # 正常中文 / 中英混排 / markdown 表格 → 不误伤
    assert not pn._chunk_text_gibberish("一、审核流程：核对发票信息并签字确认。" * 5)
    assert not pn._chunk_text_gibberish("| 名称 | 规格 | 数量 |\n| 刀叉 | A-100 | 200 |" * 5)
    assert not pn._chunk_text_gibberish("Step 1: open ERP -> click 采购订单 (PO-2026-001)" * 5)
    # PUA 乱码块 / latin-1 mojibake 块 → 拦截
    assert pn._chunk_text_gibberish("".join(chr(0xE000 + i % 200) for i in range(100)))
    assert pn._chunk_text_gibberish("ÀÂÃÄÅÆÇÈÉÊËÌÍÎÏÐÑÒÓÔÕÖ×ØÙÚÛÜÝÞß" * 5)
    # 短文本不判
    assert not pn._chunk_text_gibberish("×÷±")


def test_g25_validator_drops_gibberish_chunk():
    import opensearch_pipeline.pipeline_nodes as pn
    from opensearch_pipeline.chunker import Chunk
    good = Chunk(chunk_id="ok", doc_id="D", version_no=1, chunk_index=0,
                 chunk_type="text_chunk", chunk_text="正常的中文正文内容。" * 20, token_count=100)
    bad = Chunk(chunk_id="bad", doc_id="D", version_no=1, chunk_index=1,
                chunk_type="text_chunk",
                chunk_text="".join(chr(0xE000 + i % 200) for i in range(200)), token_count=100)
    ctx = {"chunks": [good, bad]}
    pn.node_validate_chunks(ctx)
    assert [c.chunk_id for c in ctx["valid_chunks"]] == ["ok"]
    assert ctx["invalid_chunks"][0]["chunk_id"] == "bad"
    assert "gibberish_text" in ctx["invalid_chunks"][0]["issues"]


# ═══════════════════ G12: 页眉页脚数字归一 + 非正立词过滤 ═══════════════════

def _make_pdf(tmp_path, page_builder, n_pages=3):
    import fitz
    doc = fitz.open()
    for i in range(n_pages):
        page_builder(doc.new_page(), i)
    path = str(tmp_path / "g12.pdf")
    doc.save(path)
    doc.close()
    return path


def test_g12_varying_page_number_footer_is_cropped(tmp_path):
    """"Page 1/Page 2/Page 3" 逐页变化的页码：精确文本键下永达不到频次阈值，
    数字归一（"Page #"）后按模式聚合 → 页脚被几何裁剪。"""
    from opensearch_pipeline.extraction.pdf_extractor import extract_pdf

    def build(page, i):
        page.insert_text((72, 300), "Body alpha beta gamma delta " * 4)
        page.insert_text((280, 770), f"Page {i + 1}")

    blocks, page_count, warnings = extract_pdf(_make_pdf(tmp_path, build))
    text = "\n".join(b.text for b in blocks)
    assert "Body alpha" in text
    assert "Page 1" not in text and "Page 2" not in text


def test_g12_rotated_watermark_words_dropped(tmp_path):
    """非正立（旋转）文字——典型为斜置水印——不得混入正文行。"""
    from opensearch_pipeline.extraction.pdf_extractor import extract_pdf

    def build(page, i):
        page.insert_text((72, 300), "Normal body sentence one " * 4)
        page.insert_text((300, 500), "WATERMARKTEXT", rotate=90)

    blocks, page_count, warnings = extract_pdf(_make_pdf(tmp_path, build, n_pages=1))
    text = "\n".join(b.text for b in blocks)
    assert "Normal body sentence" in text
    assert "WATERMARKTEXT" not in text


def test_g4_warning_text_emitted():
    from opensearch_pipeline.extraction.unified_extractor import _warn_skipped_vector_images
    warnings = []
    _warn_skipped_vector_images(warnings, {"skipped_vector_images": 3,
                                           "skipped_vector_formats": {"emf", "svg"}})
    assert len(warnings) == 1
    assert "[VECTOR_IMAGE_SKIPPED]" in warnings[0]
    assert "3 张" in warnings[0] and "emf/svg" in warnings[0]
    # 零跳过 → 不产生 warning
    _warn_skipped_vector_images(warnings, {})
    assert len(warnings) == 1
