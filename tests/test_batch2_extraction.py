# -*- coding: utf-8 -*-
"""Batch 2 extraction fixes — regression tests.

B2-1 xlsx hidden/veryHidden sheets are NOT ingested (text path).
B2-2 PPTX speaker notes are default-OFF (opt-in via RAG_PPTX_INCLUDE_NOTES).
B2-4 xlsx date cells render clean (no '00:00:00' noise), without needing number_format.
"""
import os
from datetime import date

import pytest

from opensearch_pipeline.extraction.unified_extractor import UnifiedExtractor


def _task(local_path, file_ext, doc_id="D1"):
    return {
        "doc_id": doc_id, "version_no": 1, "file_ext": file_ext,
        "raw_key": f"raw/test/{os.path.basename(local_path)}",
        "local_path": local_path, "filename": os.path.basename(local_path),
    }


# ── B2-1 + B2-4 (xlsx) ───────────────────────────────────────────────────────

def _make_xlsx(path, *, with_hidden=True, with_date=False):
    openpyxl = pytest.importorskip("openpyxl")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "可见表"
    ws["A1"] = "部位"
    ws["B1"] = "清扫方法"
    ws["A2"] = "VISIBLEMARKER_设备本体"
    ws["B2"] = "每班用抹布擦拭一次"
    if with_date:
        ws["C1"] = "生效日期"
        ws["C2"] = date(2024, 1, 15)
    if with_hidden:
        hs = wb.create_sheet("隐藏表")
        hs["A1"] = "HIDDENMARKER_不应入库"
        hs["A2"] = "VERYSECRET_内部源数据"
        hs.sheet_state = "hidden"
    wb.save(path)


def test_b2_1_hidden_sheet_not_ingested(tmp_path):
    p = str(tmp_path / "h.xlsx")
    _make_xlsx(p, with_hidden=True)
    res = UnifiedExtractor(simulate=True).extract(_task(p, "xlsx"))
    assert "VISIBLEMARKER_设备本体" in res.text, "可见 sheet 内容应入库"
    assert "HIDDENMARKER_不应入库" not in res.text, "隐藏 sheet 内容不应入库"
    assert "VERYSECRET_内部源数据" not in res.text
    # 不静默：跳过留痕在 warnings
    assert any("隐藏表" in w for w in res.warnings), f"隐藏 sheet 跳过应有 warning: {res.warnings}"


def test_b2_4_date_cell_renders_clean(tmp_path):
    p = str(tmp_path / "d.xlsx")
    _make_xlsx(p, with_hidden=False, with_date=True)
    res = UnifiedExtractor(simulate=True).extract(_task(p, "xlsx"))
    assert "2024-01-15" in res.text, f"日期应出现: {res.text!r}"
    assert "2024-01-15 00:00:00" not in res.text, "日期午夜时间噪声应被清理"


# ── B2-2 (pptx notes) ────────────────────────────────────────────────────────

def _make_pptx(path):
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank-ish w/ title
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tb.text_frame.text = "SLIDEBODY_采购流程说明"
    slide.notes_slide.notes_text_frame.text = "INTERNALNOTE_此页不要展示_仅讲者参考"
    prs.save(path)


def test_b2_2_speaker_notes_off_by_default(tmp_path):
    p = str(tmp_path / "deck.pptx")
    _make_pptx(p)
    res = UnifiedExtractor(simulate=True).extract(_task(p, "pptx"))
    assert "SLIDEBODY_采购流程说明" in res.text, "幻灯片正文应入库"
    assert "INTERNALNOTE_此页不要展示_仅讲者参考" not in res.text, "演讲备注默认不入库"
    assert not any(getattr(b, "source", "") == "speaker_notes" for b in res.blocks)


def test_b2_2_speaker_notes_opt_in(tmp_path, monkeypatch):
    p = str(tmp_path / "deck2.pptx")
    _make_pptx(p)
    monkeypatch.setenv("RAG_PPTX_INCLUDE_NOTES", "1")
    res = UnifiedExtractor(simulate=True).extract(_task(p, "pptx"))
    assert "INTERNALNOTE_此页不要展示_仅讲者参考" in res.text, "显式开启后备注应入库"
    assert any(getattr(b, "source", "") == "speaker_notes" for b in res.blocks)
